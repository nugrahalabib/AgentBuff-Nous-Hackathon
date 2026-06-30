"""media_providers.py — Bridge-side universal vision/video/document chains.

INTENTIONAL DUPLICATE of `hermes_plugin_files/__init__.py` chain logic.
Reason: the plugin's `__init__.py` has side effects at import time
(monkey-patches `tools.transcription_tools` + `tools.vision_tools`),
which is fine when Hermes' plugin loader runs it but inappropriate for
the bridge process. So we keep two copies of the pure chain logic:

  * Plugin copy:  hermes_plugin_files/__init__.py  — used by channels
                  (Telegram/WhatsApp/Discord/Slack) via patched Hermes
                  functions + pre_gateway_dispatch hook.
  * Bridge copy:  this file                        — used by /app via
                  attachment_preprocessor.handle_chat_send.

Both copies port OpenClaw's media-understanding architecture
(`Reff/.archive-openclaw-2026-05-21/src/media-understanding/`).
KEEP IN SYNC — when you add a provider or change a wire shape, update
BOTH files. The integration test `test_media_providers_parity()` in
scripts/ diffs the provider matrices.

What this module provides
=========================

  * transcribe_image_via_chain(image_bytes, mime, prompt=None)
       → (description, provider_used, attempts) — full vision chain
       with active-chat-model fallback.

  * transcribe_video_via_chain(video_bytes, mime, prompt=None)
       → (description, provider_used, attempts) — video chain.

  * extract_document_text(file_path, mime)
       → (text, doc_kind, error) — PDF/DOCX/XLSX/PPTX text extraction.

  * active_supports_native_pdf() → bool — True if user's active chat
       model is Anthropic Claude or Gemini (both ingest PDF inline).

  * list_configured_image_providers() → list[str] — which vision keys
       are set (for diagnostics).

Zero Hermes imports — reads `~/.hermes/config.yaml` + `~/.hermes/.env`
directly via minimal parsers (no `hermes_cli.config` dependency).
Survives `pip install --upgrade hermes-agent`.
"""

from __future__ import annotations

import base64
import logging
import os
import re
from pathlib import Path
from typing import Any, Callable, Optional

import httpx

log = logging.getLogger("bridge.media_providers")


# ──────────────────────────────────────────────────────────────────────
# Provider registry — keep in sync with hermes_plugin_files/__init__.py
# ──────────────────────────────────────────────────────────────────────


DEFAULT_MODELS_BY_CAPABILITY: dict[str, dict[str, str]] = {
    "image": {
        "openai": "gpt-5.4-mini",
        "openai-codex": "gpt-5.4",
        "anthropic": "claude-opus-4-7",
        "gemini": "gemini-3-flash-preview",
        "google": "gemini-3-flash-preview",
        "openrouter": "openrouter/auto",
        "minimax": "MiniMax-VL-01",
        "minimax-portal": "MiniMax-VL-01",
        "zai": "glm-4.6v",
        "qwen": "qwen-vl-max-latest",
        "moonshot": "kimi-k2.5",
        "xai": "grok-2-vision-latest",
    },
    "video": {
        "gemini": "gemini-3-flash-preview",
        "google": "gemini-3-flash-preview",
        "qwen": "qwen-vl-max-latest",
        "moonshot": "kimi-k2.5",
    },
}


AUTO_PRIORITY_BY_CAPABILITY: dict[str, dict[str, int]] = {
    "image": {
        "openai": 10,
        "anthropic": 20,
        "gemini": 30,
        "google": 30,
        "minimax": 40,
        "qwen": 45,
        "moonshot": 50,
        "minimax-portal": 55,
        "zai": 60,
        "openrouter": 70,
        "xai": 80,
        # openai-codex intentionally NO priority — only via active fallback
    },
    "video": {
        "gemini": 10,
        "google": 10,
        "qwen": 15,
        "moonshot": 20,
    },
}


PROVIDER_BASE_URLS: dict[str, str] = {
    "openai": "https://api.openai.com/v1",
    "openai-codex": "https://api.openai.com/v1",
    "anthropic": "https://api.anthropic.com/v1",
    "gemini": "https://generativelanguage.googleapis.com/v1beta",
    "google": "https://generativelanguage.googleapis.com/v1beta",
    "minimax": "https://api.minimax.io",
    "minimax-portal": "https://api.minimax.io",
    "zai": "https://api.z.ai/api/paas/v4",
    "qwen": "https://dashscope.aliyuncs.com/compatible-mode/v1",
    "moonshot": "https://api.moonshot.cn/v1",
    "openrouter": "https://openrouter.ai/api/v1",
    "xai": "https://api.x.ai/v1",
}


PROVIDER_KEY_ENV_VARS: dict[str, tuple[str, ...]] = {
    "openai": ("OPENAI_API_KEY",),
    "openai-codex": ("OPENAI_CODEX_API_KEY", "OPENAI_API_KEY"),
    "anthropic": ("ANTHROPIC_API_KEY",),
    "gemini": ("GEMINI_API_KEY", "GOOGLE_API_KEY", "HERMES_DEFAULT_GEMINI_KEY"),
    "google": ("GOOGLE_API_KEY", "GEMINI_API_KEY", "HERMES_DEFAULT_GEMINI_KEY"),
    "minimax": ("MINIMAX_API_KEY",),
    "minimax-portal": ("MINIMAX_PORTAL_API_KEY", "MINIMAX_API_KEY"),
    "zai": ("ZAI_API_KEY", "Z_AI_API_KEY"),
    "qwen": ("QWEN_API_KEY", "DASHSCOPE_API_KEY"),
    "moonshot": ("MOONSHOT_API_KEY", "KIMI_API_KEY"),
    "openrouter": ("OPENROUTER_API_KEY",),
    "xai": ("XAI_API_KEY", "GROK_API_KEY"),
}


_ACTIVE_PROVIDER_ALIASES = {
    "google": "google", "gemini": "gemini", "openai": "openai",
    "openai-codex": "openai-codex", "codex": "openai-codex", "codex-cli": "openai-codex",
    "anthropic": "anthropic", "claude": "anthropic", "openrouter": "openrouter",
    "kilocode": "openai", "ai-gateway": "openai",
    "minimax": "minimax", "minimax-portal": "minimax-portal",
    "zai": "zai", "z-ai": "zai", "glm": "zai",
    "qwen": "qwen", "dashscope": "qwen",
    "moonshot": "moonshot", "kimi": "moonshot",
    "xai": "xai", "grok": "xai",
}


_IMAGE_MIME_MAP = {
    "image/jpg": "image/jpeg", "image/jpeg": "image/jpeg",
    "image/png": "image/png", "image/gif": "image/gif",
    "image/webp": "image/webp", "image/bmp": "image/bmp",
}


_VIDEO_MIME_MAP = {
    "video/mp4": "video/mp4", "video/mpeg": "video/mpeg",
    "video/quicktime": "video/quicktime", "video/mov": "video/quicktime",
    "video/webm": "video/webm",
    "video/x-msvideo": "video/mp4", "video/avi": "video/mp4",
}


# ──────────────────────────────────────────────────────────────────────
# Env + config resolution (zero Hermes imports)
# ──────────────────────────────────────────────────────────────────────


def _load_dot_env_file() -> dict[str, str]:
    cached = getattr(_load_dot_env_file, "_cache", None)
    if cached is not None:
        return cached  # type: ignore[return-value]
    out: dict[str, str] = {}
    home = os.environ.get("HERMES_HOME") or str(Path.home() / ".hermes")
    env_path = Path(home) / ".env"
    try:
        if env_path.is_file():
            for raw in env_path.read_text(
                encoding="utf-8", errors="replace"
            ).splitlines():
                line = raw.strip()
                if not line or line.startswith("#") or "=" not in line:
                    continue
                key, _, value = line.partition("=")
                key = key.strip()
                if key.startswith("export "):
                    key = key[len("export "):].strip()
                if not key:
                    continue
                value = value.strip()
                if (
                    len(value) >= 2
                    and value[0] == value[-1]
                    and value[0] in ("'", '"')
                ):
                    value = value[1:-1]
                out[key] = value
    except Exception as exc:
        log.debug("media_providers: read .env failed: %s", exc)
    setattr(_load_dot_env_file, "_cache", out)
    return out


def _resolve_env_value(*var_names: str) -> Optional[str]:
    for name in var_names:
        v = os.environ.get(name)
        if v and v.strip():
            return v.strip()
    dot_env = _load_dot_env_file()
    for name in var_names:
        v = dot_env.get(name)
        if v and v.strip():
            return v.strip()
    return None


def _get_provider_api_key(provider_id: str) -> Optional[str]:
    var_names = PROVIDER_KEY_ENV_VARS.get(provider_id.lower())
    if not var_names:
        return None
    return _resolve_env_value(*var_names)


def _get_active_chat_provider() -> Optional[str]:
    home = os.environ.get("HERMES_HOME") or str(Path.home() / ".hermes")
    cfg_path = Path(home) / "config.yaml"
    if not cfg_path.is_file():
        return None
    try:
        text = cfg_path.read_text(encoding="utf-8", errors="replace")
    except Exception:
        return None
    match = re.search(r"^model\s*:\s*\n((?:\s+.+\n?)+)", text, flags=re.MULTILINE)
    if match:
        block = match.group(1)
        sub = re.search(
            r"^\s+(?:default|model)\s*:\s*['\"]?([^'\"\n#]+)['\"]?",
            block, flags=re.MULTILINE,
        )
        candidate = sub.group(1).strip() if sub else None
    else:
        inline = re.search(
            r"^model\s*:\s*['\"]?([^'\"\n#]+)['\"]?", text, flags=re.MULTILINE
        )
        candidate = inline.group(1).strip() if inline else None
    if not candidate or "/" not in candidate:
        return None
    prefix = candidate.split("/", 1)[0].strip().lower()
    return _ACTIVE_PROVIDER_ALIASES.get(prefix)


def active_supports_native_pdf() -> bool:
    """True if active chat model can ingest PDF inline (Anthropic/Gemini)."""
    return _get_active_chat_provider() in {"anthropic", "gemini", "google"}


def list_configured_image_providers() -> list[str]:
    """Diagnostic: which vision providers have an API key configured."""
    out: list[str] = []
    seen: set[str] = set()
    for provider_id in PROVIDER_KEY_ENV_VARS.keys():
        canonical = "gemini" if provider_id == "google" else provider_id
        if canonical in seen:
            continue
        if _get_provider_api_key(provider_id):
            seen.add(canonical)
            out.append(canonical)
    return sorted(out)


def _read_yaml_tools_media_value(*keys: str) -> Optional[Any]:
    """Read a value from ~/.hermes/config.yaml::tools.media.<keys>."""
    home = os.environ.get("HERMES_HOME") or str(Path.home() / ".hermes")
    cfg_path = Path(home) / "config.yaml"
    if not cfg_path.is_file():
        return None
    try:
        import yaml
        with open(cfg_path, "r", encoding="utf-8") as f:
            data = yaml.safe_load(f) or {}
    except Exception:
        return None
    if not isinstance(data, dict):
        return None
    node: Any = data.get("tools")
    if not isinstance(node, dict):
        return None
    node = node.get("media")
    if not isinstance(node, dict):
        return None
    for key in keys:
        if not isinstance(node, dict):
            return None
        node = node.get(key)
        if node is None:
            return None
    return node


def _capability_enabled(capability: str) -> bool:
    """True unless tools.media.<capability>.enabled is explicitly false."""
    value = _read_yaml_tools_media_value(capability, "enabled")
    return value is not False


def _read_provider_overrides(provider_id: str) -> dict:
    """Per-provider overrides from tools.media.providers.<id>.* in config.yaml."""
    overrides = _read_yaml_tools_media_value("providers", provider_id)
    return overrides if isinstance(overrides, dict) else {}


# ─────────────────────────────────────────────────────────────────────
# Binary cache — LRU file-bytes cache by (path, mtime, size)
# ─────────────────────────────────────────────────────────────────────


class _BinaryCache:
    def __init__(self, max_entries: int = 32) -> None:
        import collections
        self._max = max_entries
        self._bytes: "collections.OrderedDict[tuple, bytes]" = collections.OrderedDict()

    def _key(self, file_path: str) -> Optional[tuple]:
        try:
            st = Path(file_path).stat()
            return (str(Path(file_path).resolve()), st.st_mtime_ns, st.st_size)
        except Exception:
            return None

    def get(self, file_path: str) -> Optional[bytes]:
        key = self._key(file_path)
        if key is None or key not in self._bytes:
            return None
        self._bytes.move_to_end(key)
        return self._bytes[key]

    def put(self, file_path: str, data: bytes) -> None:
        key = self._key(file_path)
        if key is None:
            return
        self._bytes[key] = data
        self._bytes.move_to_end(key)
        while len(self._bytes) > self._max:
            self._bytes.popitem(last=False)


_BINARY_CACHE = _BinaryCache()


def _read_file_cached(file_path: str) -> Optional[bytes]:
    cached = _BINARY_CACHE.get(file_path)
    if cached is not None:
        return cached
    try:
        data = Path(file_path).read_bytes()
    except Exception:
        return None
    _BINARY_CACHE.put(file_path, data)
    return data


# ─────────────────────────────────────────────────────────────────────
# Concurrent execution
# ─────────────────────────────────────────────────────────────────────


def _media_concurrency_limit() -> int:
    raw = os.environ.get("MEDIA_CONCURRENCY") or os.environ.get("HERMES_MEDIA_CONCURRENCY")
    if not raw:
        return 3
    try:
        return max(1, min(int(raw), 16))
    except (TypeError, ValueError):
        return 3


def run_concurrently(tasks: list) -> list:
    """Run callables in a thread pool. Exceptions become exception objects in result list."""
    if not tasks:
        return []
    if len(tasks) == 1:
        try:
            return [tasks[0]()]
        except Exception as exc:
            return [exc]
    from concurrent.futures import ThreadPoolExecutor
    limit = min(len(tasks), _media_concurrency_limit())
    with ThreadPoolExecutor(max_workers=limit) as pool:
        futures = [pool.submit(t) for t in tasks]
        return [
            (f.result() if not f.exception() else f.exception())
            for f in futures
        ]


def _http_post(
    url: str, *, provider_id: Optional[str] = None, **kwargs
) -> httpx.Response:
    """HTTP POST with per-provider TLS/proxy/authHeader override support."""
    timeout = httpx.Timeout(connect=10.0, read=90.0, write=90.0, pool=10.0)
    client_kwargs: dict[str, Any] = {"timeout": timeout}
    if provider_id:
        overrides = _read_provider_overrides(provider_id)
        if overrides:
            tls_verify = overrides.get("tlsVerify")
            if tls_verify is False:
                client_kwargs["verify"] = False
            elif tls_verify is True:
                client_kwargs["verify"] = True
            tls_ca = overrides.get("tlsCaCert")
            if isinstance(tls_ca, str) and tls_ca.strip():
                client_kwargs["verify"] = tls_ca.strip()
            proxy_url = overrides.get("proxy")
            if isinstance(proxy_url, str) and proxy_url.strip():
                client_kwargs["proxy"] = proxy_url.strip()
            auth = overrides.get("authHeader")
            if isinstance(auth, dict):
                name = auth.get("name")
                value = auth.get("value")
                if isinstance(name, str) and isinstance(value, str):
                    headers = kwargs.setdefault("headers", {})
                    if isinstance(headers, dict):
                        headers[name] = value
    with httpx.Client(**client_kwargs) as client:
        return client.post(url, **kwargs)


# ──────────────────────────────────────────────────────────────────────
# VISION provider implementations
# ──────────────────────────────────────────────────────────────────────


def _describe_openai_compatible_image(
    provider_id: str,
    image_bytes: bytes,
    mime: str,
    prompt: str,
) -> tuple[Optional[str], Optional[str]]:
    api_key = _get_provider_api_key(provider_id)
    if not api_key:
        return (None, "no_api_key")
    base_url = (
        os.environ.get(f"{provider_id.upper()}_BASE_URL")
        or PROVIDER_BASE_URLS.get(provider_id)
    )
    if not base_url:
        return (None, "no base_url")
    model = (
        os.environ.get(f"{provider_id.upper()}_VISION_MODEL")
        or DEFAULT_MODELS_BY_CAPABILITY["image"].get(provider_id)
    )
    if not model:
        return (None, "no default vision model")
    canon_mime = _IMAGE_MIME_MAP.get((mime or "").lower(), "image/jpeg")
    try:
        b64 = base64.b64encode(image_bytes).decode("ascii")
        body = {
            "model": model,
            "messages": [
                {
                    "role": "user",
                    "content": [
                        {"type": "text", "text": prompt},
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:{canon_mime};base64,{b64}",
                            },
                        },
                    ],
                }
            ],
            "max_tokens": 1024,
        }
        resp = _http_post(
            f"{base_url.rstrip('/')}/chat/completions",
            provider_id=provider_id,
            headers={
                "Authorization": f"Bearer {api_key}",
                "Content-Type": "application/json",
            },
            json=body,
        )
        if resp.status_code != 200:
            return (None, f"HTTP {resp.status_code}: {resp.text[:200]}")
        data = resp.json()
        choices = data.get("choices") or []
        if not choices:
            return (None, "no choices")
        message = choices[0].get("message") or {}
        content = message.get("content")
        if isinstance(content, str):
            return (content.strip(), None)
        if isinstance(content, list):
            text = "\n".join(
                (p.get("text") or "").strip()
                for p in content
                if isinstance(p, dict) and p.get("text")
            ).strip()
            if text:
                return (text, None)
        reasoning = message.get("reasoning_content")
        if isinstance(reasoning, str) and reasoning.strip():
            return (reasoning.strip(), None)
        return (None, "empty content")
    except Exception as exc:
        return (None, f"{type(exc).__name__}: {exc}")


def _describe_anthropic_image(
    image_bytes: bytes, mime: str, prompt: str
) -> tuple[Optional[str], Optional[str]]:
    api_key = _get_provider_api_key("anthropic")
    if not api_key:
        return (None, "no_api_key")
    model = (
        os.environ.get("ANTHROPIC_VISION_MODEL")
        or DEFAULT_MODELS_BY_CAPABILITY["image"]["anthropic"]
    )
    canon_mime = _IMAGE_MIME_MAP.get((mime or "").lower(), "image/jpeg")
    try:
        b64 = base64.b64encode(image_bytes).decode("ascii")
        body = {
            "model": model,
            "max_tokens": 1024,
            "messages": [
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "image",
                            "source": {
                                "type": "base64",
                                "media_type": canon_mime,
                                "data": b64,
                            },
                        },
                        {"type": "text", "text": prompt},
                    ],
                }
            ],
        }
        resp = _http_post(
            f"{PROVIDER_BASE_URLS['anthropic']}/messages",
            provider_id="anthropic",
            headers={
                "x-api-key": api_key,
                "anthropic-version": "2023-06-01",
                "content-type": "application/json",
            },
            json=body,
        )
        if resp.status_code != 200:
            return (None, f"HTTP {resp.status_code}: {resp.text[:200]}")
        data = resp.json()
        content = data.get("content") or []
        text = "\n".join(
            (p.get("text") or "").strip()
            for p in content
            if isinstance(p, dict) and p.get("type") == "text"
        ).strip()
        return ((text or None), None if text else "empty content")
    except Exception as exc:
        return (None, f"{type(exc).__name__}: {exc}")


def _describe_gemini_image(
    image_bytes: bytes, mime: str, prompt: str
) -> tuple[Optional[str], Optional[str]]:
    api_key = _get_provider_api_key("gemini")
    if not api_key:
        return (None, "no_api_key")
    model = (
        os.environ.get("GEMINI_VISION_MODEL")
        or DEFAULT_MODELS_BY_CAPABILITY["image"]["gemini"]
    )
    canon_mime = _IMAGE_MIME_MAP.get((mime or "").lower(), "image/jpeg")
    try:
        b64 = base64.b64encode(image_bytes).decode("ascii")
        body = {
            "contents": [{
                "parts": [
                    {"text": prompt},
                    {"inline_data": {"mime_type": canon_mime, "data": b64}},
                ]
            }],
            "generationConfig": {"temperature": 0.1, "maxOutputTokens": 1024},
        }
        url = (
            f"{PROVIDER_BASE_URLS['gemini']}/models/{model}:generateContent"
            f"?key={api_key}"
        )
        resp = _http_post(url, provider_id="gemini", json=body)
        if resp.status_code != 200:
            return (None, f"HTTP {resp.status_code}: {resp.text[:200]}")
        data = resp.json()
        candidates = data.get("candidates") or []
        if not candidates:
            return (None, "no candidates")
        parts = (candidates[0].get("content") or {}).get("parts") or []
        text = "".join(
            (p.get("text") or "") for p in parts if isinstance(p, dict)
        ).strip()
        return ((text or None), None if text else "empty content")
    except Exception as exc:
        return (None, f"{type(exc).__name__}: {exc}")


def _describe_minimax_image(
    image_bytes: bytes, mime: str, prompt: str
) -> tuple[Optional[str], Optional[str]]:
    api_key = _get_provider_api_key("minimax")
    if not api_key:
        return (None, "no_api_key")
    canon_mime = _IMAGE_MIME_MAP.get((mime or "").lower(), "image/jpeg")
    base_url = (
        os.environ.get("MINIMAX_BASE_URL")
        or PROVIDER_BASE_URLS["minimax"]
    )
    try:
        b64 = base64.b64encode(image_bytes).decode("ascii")
        body = {
            "prompt": prompt,
            "image_url": f"data:{canon_mime};base64,{b64}",
        }
        resp = _http_post(
            f"{base_url.rstrip('/')}/v1/coding_plan/vlm",
            provider_id="minimax",
            headers={
                "Authorization": f"Bearer {api_key}",
                "Content-Type": "application/json",
                "MM-API-Source": "AgentBuff",
            },
            json=body,
        )
        if resp.status_code != 200:
            return (None, f"HTTP {resp.status_code}: {resp.text[:200]}")
        data = resp.json()
        base_resp = data.get("base_resp") or {}
        status_code = base_resp.get("status_code")
        if status_code not in (None, 0):
            return (None, f"minimax error {status_code}: {base_resp.get('status_msg')}")
        content = data.get("content")
        if isinstance(content, str) and content.strip():
            return (content.strip(), None)
        return (None, "empty content")
    except Exception as exc:
        return (None, f"{type(exc).__name__}: {exc}")


def _describe_minimax_image_with_provider(
    provider_id: str, image_bytes: bytes, mime: str, prompt: str
) -> tuple[Optional[str], Optional[str]]:
    """Variant of _describe_minimax_image that picks the API key per
    provider_id — distinguishes minimax vs minimax-portal."""
    api_key = _get_provider_api_key(provider_id)
    if not api_key:
        return (None, "no_api_key")
    canon_mime = _IMAGE_MIME_MAP.get((mime or "").lower(), "image/jpeg")
    base_url = (
        os.environ.get(f"{provider_id.upper().replace('-', '_')}_BASE_URL")
        or PROVIDER_BASE_URLS.get(provider_id)
    )
    if not base_url:
        return (None, "no base_url")
    try:
        b64 = base64.b64encode(image_bytes).decode("ascii")
        body = {"prompt": prompt, "image_url": f"data:{canon_mime};base64,{b64}"}
        resp = _http_post(
            f"{base_url.rstrip('/')}/v1/coding_plan/vlm",
            provider_id=provider_id,
            headers={
                "Authorization": f"Bearer {api_key}",
                "Content-Type": "application/json",
                "MM-API-Source": "AgentBuff",
            },
            json=body,
        )
        if resp.status_code != 200:
            return (None, f"HTTP {resp.status_code}: {resp.text[:200]}")
        data = resp.json()
        base_resp = data.get("base_resp") or {}
        status_code = base_resp.get("status_code")
        if status_code not in (None, 0):
            return (None, f"{provider_id} error {status_code}: {base_resp.get('status_msg')}")
        content = data.get("content")
        if isinstance(content, str) and content.strip():
            return (content.strip(), None)
        return (None, "empty content")
    except Exception as exc:
        return (None, f"{type(exc).__name__}: {exc}")


PROVIDER_DESCRIBE_IMAGE_FNS: dict[
    str, Callable[[bytes, str, str], tuple[Optional[str], Optional[str]]]
] = {
    "openai": lambda b, m, p: _describe_openai_compatible_image("openai", b, m, p),
    "openai-codex": lambda b, m, p: _describe_openai_compatible_image("openai-codex", b, m, p),
    "openrouter": lambda b, m, p: _describe_openai_compatible_image("openrouter", b, m, p),
    "xai": lambda b, m, p: _describe_openai_compatible_image("xai", b, m, p),
    "qwen": lambda b, m, p: _describe_openai_compatible_image("qwen", b, m, p),
    "moonshot": lambda b, m, p: _describe_openai_compatible_image("moonshot", b, m, p),
    "zai": lambda b, m, p: _describe_openai_compatible_image("zai", b, m, p),
    "anthropic": _describe_anthropic_image,
    "gemini": _describe_gemini_image,
    "google": _describe_gemini_image,
    "minimax": _describe_minimax_image,
    "minimax-portal": lambda b, m, p: _describe_minimax_image_with_provider(
        "minimax-portal", b, m, p
    ),
}


def _default_image_prompt(user_prompt: Optional[str]) -> str:
    if user_prompt and user_prompt.strip():
        return user_prompt.strip()
    return (
        "Describe this image in detail. Note any text visible in the image "
        "verbatim. If it's a document or screenshot, transcribe the content. "
        "If it's a chart, describe data trends. Be thorough but concise."
    )


def transcribe_image_via_chain(
    image_bytes: bytes,
    mime: str,
    prompt: Optional[str] = None,
) -> tuple[Optional[str], Optional[str], list[str]]:
    """Run vision chain: active → bundled priority order.

    Returns (description, provider_used, attempts). description=None
    means no provider produced output.
    """
    if not _capability_enabled("image"):
        return (None, None, ["image capability disabled in config"])
    canon_mime = _IMAGE_MIME_MAP.get((mime or "").lower(), "image/jpeg")
    actual_prompt = _default_image_prompt(prompt)
    attempts: list[str] = []
    tried: set[str] = set()

    def _canon(pid: str) -> str:
        return "gemini" if pid == "google" else pid

    def _try(provider_id: str) -> Optional[str]:
        canonical = _canon(provider_id)
        if canonical in tried:
            return None
        tried.add(canonical)
        fn = PROVIDER_DESCRIBE_IMAGE_FNS.get(provider_id)
        if not fn:
            attempts.append(f"{provider_id}: not registered")
            return None
        text, err = fn(image_bytes, canon_mime, actual_prompt)
        if text:
            attempts.append(f"{provider_id}: ok ({len(text)} chars)")
            return text
        attempts.append(f"{provider_id}: {err or 'empty'}")
        return None

    active = _get_active_chat_provider()
    if active:
        result = _try(active)
        if result:
            return (result, active, attempts)

    items = sorted(
        AUTO_PRIORITY_BY_CAPABILITY["image"].items(),
        key=lambda kv: (kv[1], kv[0]),
    )
    for provider_id, _prio in items:
        result = _try(provider_id)
        if result:
            return (result, provider_id, attempts)

    return (None, None, attempts)


# ──────────────────────────────────────────────────────────────────────
# VIDEO provider implementations
# ──────────────────────────────────────────────────────────────────────


def _describe_video_openai_compat(
    provider_id: str,
    video_bytes: bytes,
    mime: str,
    prompt: str,
) -> tuple[Optional[str], Optional[str]]:
    api_key = _get_provider_api_key(provider_id)
    if not api_key:
        return (None, "no_api_key")
    base_url = (
        os.environ.get(f"{provider_id.upper()}_BASE_URL")
        or PROVIDER_BASE_URLS.get(provider_id)
    )
    if not base_url:
        return (None, "no base_url")
    model = (
        os.environ.get(f"{provider_id.upper()}_VIDEO_MODEL")
        or DEFAULT_MODELS_BY_CAPABILITY["video"].get(provider_id)
    )
    if not model:
        return (None, "no default video model")
    canon_mime = _VIDEO_MIME_MAP.get((mime or "").lower(), "video/mp4")
    try:
        b64 = base64.b64encode(video_bytes).decode("ascii")
        body = {
            "model": model,
            "messages": [{
                "role": "user",
                "content": [
                    {"type": "text", "text": prompt},
                    {
                        "type": "video_url",
                        "video_url": {
                            "url": f"data:{canon_mime};base64,{b64}",
                        },
                    },
                ],
            }],
            "max_tokens": 1024,
        }
        resp = _http_post(
            f"{base_url.rstrip('/')}/chat/completions",
            provider_id=provider_id,
            headers={
                "Authorization": f"Bearer {api_key}",
                "Content-Type": "application/json",
            },
            json=body,
        )
        if resp.status_code != 200:
            return (None, f"HTTP {resp.status_code}: {resp.text[:200]}")
        data = resp.json()
        choices = data.get("choices") or []
        if not choices:
            return (None, "no choices")
        message = choices[0].get("message") or {}
        content = message.get("content")
        if isinstance(content, str):
            return (content.strip(), None)
        if isinstance(content, list):
            text = "\n".join(
                (p.get("text") or "").strip()
                for p in content
                if isinstance(p, dict) and p.get("text")
            ).strip()
            if text:
                return (text, None)
        return (None, "empty content")
    except Exception as exc:
        return (None, f"{type(exc).__name__}: {exc}")


def _describe_gemini_video(
    video_bytes: bytes, mime: str, prompt: str
) -> tuple[Optional[str], Optional[str]]:
    api_key = _get_provider_api_key("gemini")
    if not api_key:
        return (None, "no_api_key")
    model = (
        os.environ.get("GEMINI_VIDEO_MODEL")
        or DEFAULT_MODELS_BY_CAPABILITY["video"]["gemini"]
    )
    canon_mime = _VIDEO_MIME_MAP.get((mime or "").lower(), "video/mp4")
    try:
        b64 = base64.b64encode(video_bytes).decode("ascii")
        body = {
            "contents": [{
                "parts": [
                    {"text": prompt},
                    {"inline_data": {"mime_type": canon_mime, "data": b64}},
                ]
            }],
            "generationConfig": {"temperature": 0.1, "maxOutputTokens": 1024},
        }
        url = (
            f"{PROVIDER_BASE_URLS['gemini']}/models/{model}:generateContent"
            f"?key={api_key}"
        )
        resp = _http_post(url, provider_id="gemini", json=body)
        if resp.status_code != 200:
            return (None, f"HTTP {resp.status_code}: {resp.text[:200]}")
        data = resp.json()
        candidates = data.get("candidates") or []
        if not candidates:
            return (None, "no candidates")
        parts = (candidates[0].get("content") or {}).get("parts") or []
        text = "".join(
            (p.get("text") or "") for p in parts if isinstance(p, dict)
        ).strip()
        return ((text or None), None if text else "empty content")
    except Exception as exc:
        return (None, f"{type(exc).__name__}: {exc}")


PROVIDER_DESCRIBE_VIDEO_FNS: dict[
    str, Callable[[bytes, str, str], tuple[Optional[str], Optional[str]]]
] = {
    "gemini": _describe_gemini_video,
    "google": _describe_gemini_video,
    "qwen": lambda b, m, p: _describe_video_openai_compat("qwen", b, m, p),
    "moonshot": lambda b, m, p: _describe_video_openai_compat("moonshot", b, m, p),
}


def _default_video_prompt(user_prompt: Optional[str]) -> str:
    if user_prompt and user_prompt.strip():
        return user_prompt.strip()
    return (
        "Describe this video in detail. Note key moments, on-screen text, "
        "spoken audio if any, and overall content. Be thorough but concise."
    )


def transcribe_video_via_chain(
    video_bytes: bytes,
    mime: str,
    prompt: Optional[str] = None,
) -> tuple[Optional[str], Optional[str], list[str]]:
    if not _capability_enabled("video"):
        return (None, None, ["video capability disabled in config"])
    canon_mime = _VIDEO_MIME_MAP.get((mime or "").lower(), "video/mp4")
    actual_prompt = _default_video_prompt(prompt)
    attempts: list[str] = []
    tried: set[str] = set()

    def _canon(pid: str) -> str:
        return "gemini" if pid == "google" else pid

    def _try(provider_id: str) -> Optional[str]:
        canonical = _canon(provider_id)
        if canonical in tried:
            return None
        tried.add(canonical)
        fn = PROVIDER_DESCRIBE_VIDEO_FNS.get(provider_id)
        if not fn:
            attempts.append(f"{provider_id}: not registered")
            return None
        text, err = fn(video_bytes, canon_mime, actual_prompt)
        if text:
            attempts.append(f"{provider_id}: ok ({len(text)} chars)")
            return text
        attempts.append(f"{provider_id}: {err or 'empty'}")
        return None

    active = _get_active_chat_provider()
    if active:
        result = _try(active)
        if result:
            return (result, active, attempts)

    items = sorted(
        AUTO_PRIORITY_BY_CAPABILITY["video"].items(),
        key=lambda kv: (kv[1], kv[0]),
    )
    for provider_id, _prio in items:
        result = _try(provider_id)
        if result:
            return (result, provider_id, attempts)

    return (None, None, attempts)


# ──────────────────────────────────────────────────────────────────────
# DOCUMENT text extraction
# ──────────────────────────────────────────────────────────────────────


MAX_DOC_EXTRACT_CHARS = 50_000


_DOCUMENT_MIME_HINTS = {
    "application/pdf": "pdf",
    "application/x-pdf": "pdf",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "docx",
    "application/msword": "docx",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": "xlsx",
    "application/vnd.ms-excel": "xlsx",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
    "application/vnd.ms-powerpoint": "pptx",
}

_DOCUMENT_EXT_HINTS = {
    ".pdf": "pdf", ".docx": "docx", ".doc": "docx",
    ".xlsx": "xlsx", ".xls": "xlsx",
    ".pptx": "pptx", ".ppt": "pptx",
}


def _classify_document(file_path: str, mime: str) -> Optional[str]:
    canon = (mime or "").lower()
    kind = _DOCUMENT_MIME_HINTS.get(canon)
    if kind:
        return kind
    return _DOCUMENT_EXT_HINTS.get(Path(file_path).suffix.lower())


def _extract_pdf_text(file_path: str) -> tuple[Optional[str], Optional[str]]:
    try:
        import pdfplumber  # type: ignore[import]
        try:
            with pdfplumber.open(file_path) as pdf:
                pages_text: list[str] = []
                for i, page in enumerate(pdf.pages):
                    try:
                        t = page.extract_text() or ""
                    except Exception:
                        t = ""
                    if t.strip():
                        pages_text.append(f"--- Page {i + 1} ---\n{t.strip()}")
                    if sum(len(p) for p in pages_text) >= MAX_DOC_EXTRACT_CHARS:
                        break
                combined = "\n\n".join(pages_text).strip()
                if combined:
                    return (combined[:MAX_DOC_EXTRACT_CHARS], None)
                return (None, "pdf had no extractable text")
        except Exception as exc:
            return (None, f"pdfplumber: {type(exc).__name__}: {exc}")
    except ImportError:
        pass

    try:
        import pypdf  # type: ignore[import]
        try:
            reader = pypdf.PdfReader(file_path)
            pages_text = []
            for i, page in enumerate(reader.pages):
                try:
                    t = page.extract_text() or ""
                except Exception:
                    t = ""
                if t.strip():
                    pages_text.append(f"--- Page {i + 1} ---\n{t.strip()}")
                if sum(len(p) for p in pages_text) >= MAX_DOC_EXTRACT_CHARS:
                    break
            combined = "\n\n".join(pages_text).strip()
            if combined:
                return (combined[:MAX_DOC_EXTRACT_CHARS], None)
            return (None, "pdf had no extractable text")
        except Exception as exc:
            return (None, f"pypdf: {type(exc).__name__}: {exc}")
    except ImportError:
        return (None, "no PDF library available (pdfplumber/pypdf)")


def _extract_docx_text(file_path: str) -> tuple[Optional[str], Optional[str]]:
    try:
        import docx  # type: ignore[import]
    except ImportError:
        return (None, "python-docx not installed")
    try:
        doc = docx.Document(file_path)
        chunks = [p.text for p in doc.paragraphs if p.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        chunks.append(cell.text.strip())
        combined = "\n".join(chunks).strip()
        if not combined:
            return (None, "docx had no extractable text")
        return (combined[:MAX_DOC_EXTRACT_CHARS], None)
    except Exception as exc:
        return (None, f"docx: {type(exc).__name__}: {exc}")


def _extract_xlsx_text(file_path: str) -> tuple[Optional[str], Optional[str]]:
    try:
        import openpyxl  # type: ignore[import]
    except ImportError:
        return (None, "openpyxl not installed")
    try:
        wb = openpyxl.load_workbook(file_path, data_only=True, read_only=True)
        chunks: list[str] = []
        total = 0
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            chunks.append(f"--- Sheet: {sheet_name} ---")
            for row in ws.iter_rows(values_only=True):
                row_text = " | ".join(
                    "" if v is None else str(v) for v in row
                ).strip(" |")
                if row_text:
                    chunks.append(row_text)
                    total += len(row_text)
                    if total >= MAX_DOC_EXTRACT_CHARS:
                        chunks.append("...(truncated, file too large)")
                        wb.close()
                        return ("\n".join(chunks)[:MAX_DOC_EXTRACT_CHARS], None)
            chunks.append("")
        wb.close()
        combined = "\n".join(chunks).strip()
        if not combined:
            return (None, "xlsx had no extractable text")
        return (combined[:MAX_DOC_EXTRACT_CHARS], None)
    except Exception as exc:
        return (None, f"openpyxl: {type(exc).__name__}: {exc}")


def _extract_pptx_text(file_path: str) -> tuple[Optional[str], Optional[str]]:
    try:
        from pptx import Presentation  # type: ignore[import]
    except ImportError:
        return (None, "python-pptx not installed")
    try:
        prs = Presentation(file_path)
        chunks: list[str] = []
        for i, slide in enumerate(prs.slides):
            chunks.append(f"--- Slide {i + 1} ---")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    chunks.append(shape.text.strip())
            chunks.append("")
        combined = "\n".join(chunks).strip()
        if not combined:
            return (None, "pptx had no extractable text")
        return (combined[:MAX_DOC_EXTRACT_CHARS], None)
    except Exception as exc:
        return (None, f"pptx: {type(exc).__name__}: {exc}")


_DOC_EXTRACTORS: dict[
    str, Callable[[str], tuple[Optional[str], Optional[str]]]
] = {
    "pdf": _extract_pdf_text,
    "docx": _extract_docx_text,
    "xlsx": _extract_xlsx_text,
    "pptx": _extract_pptx_text,
}


def extract_document_text(
    file_path: str, mime: str
) -> tuple[Optional[str], Optional[str], Optional[str]]:
    """Extract document text. Returns (text, doc_kind, error)."""
    if not _capability_enabled("document"):
        return (None, None, "document capability disabled in config")
    doc_kind = _classify_document(file_path, mime)
    if doc_kind is None:
        return (None, None, "not a recognized document type")
    extractor = _DOC_EXTRACTORS.get(doc_kind)
    if extractor is None:
        return (None, doc_kind, f"no extractor for {doc_kind}")
    text, err = extractor(file_path)
    return (text, doc_kind, err)


__all__ = [
    "transcribe_image_via_chain",
    "transcribe_video_via_chain",
    "extract_document_text",
    "active_supports_native_pdf",
    "list_configured_image_providers",
    "run_concurrently",
    "DEFAULT_MODELS_BY_CAPABILITY",
    "AUTO_PRIORITY_BY_CAPABILITY",
    "PROVIDER_BASE_URLS",
    "PROVIDER_KEY_ENV_VARS",
]
