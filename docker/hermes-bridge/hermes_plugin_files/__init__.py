"""AgentBuff multimodal plugin — universal multimodal extensions for Hermes.

Loaded by Hermes' plugin system (`hermes_cli/plugins.py::discover_plugins`)
when the plugin is enabled in `~/.hermes/config.yaml::plugins.enabled`. The
plugin lives at `$HERMES_HOME/plugins/agentbuff-multimodal/` — i.e. inside
the user's Hermes volume, NOT in the `hermes-agent` pip package. So when
`pip install --upgrade hermes-agent` overwrites the package, this plugin
stays untouched.

What this plugin does
=====================

Mirrors OpenClaw's COMPLETE media-understanding architecture
(`Reff/.archive-openclaw-2026-05-21/src/media-understanding/`) — covers
audio, image, video, AND document inputs across every channel and the
/app web UI.

  AUDIO (STT)
  -----------
  * Universal STT registry — OpenAI / Groq / Deepgram / Gemini / Mistral
    / xAI / Anthropic with auto-priority chain and active-chat-model
    fallback (reads `model.default` from config.yaml).
  * Patches `tools.transcription_tools.transcribe_audio` (called by
    `gateway/run.py:_enrich_message_with_transcription` for every audio
    attachment).

  IMAGE (VISION)
  --------------
  * Vision provider chain — OpenAI / Anthropic / Gemini / OpenRouter /
    MiniMax / Z.AI / Qwen / Moonshot / xAI with active-chat-model
    fallback.
  * Patches `tools.vision_tools.vision_analyze_tool` (called by
    `gateway/run.py:_enrich_message_with_vision` AND directly by the
    agent when it invokes `vision_analyze`).

  VIDEO
  -----
  * Video provider chain — Gemini / Qwen / Moonshot via the OpenAI-
    compatible `chat.completions` shape with `type: "video_url"` parts.
  * Hermes' gateway does NOT auto-call `video_analyze_tool` (only
    `vision_analyze_tool` is auto-invoked). To cover channels we
    intercept incoming `MessageEvent` via the `pre_gateway_dispatch`
    hook: for every `video/*` attachment we run the chain and prepend
    a `[The user sent a video. Here's what's in it: "..."]` note to
    the message text — same shape OpenClaw used.

  DOCUMENT
  --------
  * Native passthrough — when active chat model is Anthropic Claude
    (PDF) or Gemini (PDF), we just keep the existing context note
    pointing at the cached path; the model reads PDF inline via its
    own multimodal API (matches OpenClaw's `nativeDocumentInputs`
    behaviour).
  * Text extraction fallback — for any other active model we extract
    text via pdfplumber (PDF), python-docx (DOCX), openpyxl (XLSX),
    python-pptx (PPTX) and inline the text into the message.
  * Also runs via `pre_gateway_dispatch` hook so it covers channels +
    /app uniformly.

Why two strategies (monkey-patch + hook)
========================================

Different surfaces of Hermes invoke different functions:

  * `transcribe_audio` and `vision_analyze_tool` are auto-called by
    the gateway for audio/image attachments — monkey-patching them
    extends the providers available on that auto-call path. This
    matches what users expect: chief drops a voice note in Telegram
    and the agent sees the transcript on the very first turn, no
    explicit tool-call needed.

  * `video_analyze_tool` exists but the gateway does NOT auto-call
    it. Without a hook the agent would have to manually invoke
    `video_analyze` every time, which mass-market users won't know
    to do. The `pre_gateway_dispatch` hook fixes this by running the
    chain inline and rewriting the user message before dispatch.

  * Documents have no Hermes-native tool at all. The hook is the
    only insertion point.

Timing
======

Hermes calls `discover_plugins()` at startup
(`hermes_cli/gateway.py:3715`). When this plugin's `__init__.py` runs:

  1. Import-time side effects: `_install_patches()` patches the
     transcribe_audio + vision_analyze_tool symbols on
     `tools.transcription_tools` / `tools.vision_tools`, plus walks
     `sys.modules` to rebind any top-level importers.
  2. After import, Hermes calls `register(ctx)` (the standard plugin
     entry point per `hermes_cli/plugins.py:1184-1190`) which we use
     to register the `pre_gateway_dispatch` hook for video +
     document handling.

  Critical: `gateway/run.py` does the `from tools... import` calls
  INSIDE method bodies (line 14293 for transcribe, line 14187+ for
  vision), so fresh patches are picked up automatically on each
  call without needing to restart Hermes.

Safety
======

Every patched call is wrapped in try/except. If the wrapper crashes
for any reason — bad MIME, network drop, JSON parse error, Hermes
API contract change — we fall through to Hermes' ORIGINAL function.
The plugin never breaks a working pipeline; worst case it adds zero
value, never negative value. Hook callbacks have the same defensive
posture: any exception is caught and logged, dispatch continues with
the original message.

Provider matrix (verbatim from OpenClaw `bundled-defaults.ts`)
==============================================================

  | Capability | Providers (priority order)                                       |
  | audio      | openai · groq · deepgram · gemini · mistral · xai                |
  | image      | openai · anthropic · gemini · openrouter · minimax · zai · qwen  |
  |            | · moonshot · xai                                                  |
  | video      | gemini · qwen · moonshot                                          |
  | document   | (native: anthropic+gemini)  ·  (extract: pdfplumber/docx/xlsx/pptx)|

Plus active-chat-model fallback at every tier: if user's
`model.default` provider supports the capability, try it FIRST so
the user re-uses the API key they already pay for as their chat LLM.
"""

from __future__ import annotations

import base64
import io
import logging
import os
import re
import sys
from pathlib import Path
from typing import Any, Callable, Optional

logger = logging.getLogger(__name__)

# Sentinel attribute we attach to our wrapper so subsequent plugin
# reloads (if any) don't double-wrap.
_PATCHED_SENTINEL = "_agentbuff_multimodal_patched"


# ──────────────────────────────────────────────────────────────────────
# Provider registry — ported from OpenClaw's bundled-defaults.ts +
# stt_providers.py (in the bridge sibling package). Kept self-contained
# inside the plugin so a Hermes upgrade that ships a different bridge
# layout doesn't break this code.
# ──────────────────────────────────────────────────────────────────────

# Default model per (provider, capability). Mirrors OpenClaw's
# bundled-defaults.ts; pinned to the same model IDs OpenClaw shipped
# so behaviour is byte-identical for users migrating between engines.
DEFAULT_MODELS_BY_CAPABILITY: dict[str, dict[str, str]] = {
    "audio": {
        "openai": "gpt-4o-transcribe",
        "groq": "whisper-large-v3-turbo",
        "gemini": "gemini-2.5-flash",
        "google": "gemini-2.5-flash",
        "deepgram": "nova-3",
        "mistral": "voxtral-mini-latest",
        "xai": "grok-3-stt",
    },
    "image": {
        "openai": "gpt-5.4-mini",
        "openai-codex": "gpt-5.4",
        "anthropic": "claude-opus-4-7",
        "gemini": "gemini-3-flash-preview",
        "google": "gemini-3-flash-preview",
        "openrouter": "openrouter/auto",
        "minimax": "MiniMax-VL-01",
        "minimax-portal": "MiniMax-VL-01",
        "zai": "glm-4.6v",
        "qwen": "qwen-vl-max-latest",
        "moonshot": "kimi-k2.5",
        "xai": "grok-2-vision-latest",
    },
    "video": {
        "gemini": "gemini-3-flash-preview",
        "google": "gemini-3-flash-preview",
        "qwen": "qwen-vl-max-latest",
        "moonshot": "kimi-k2.5",
    },
    # Document = same models as image (vision-capable LLMs can ingest
    # PDF/image-of-page either inline (anthropic/gemini) or via the
    # extraction-then-prompt path).
    "document": {
        "anthropic": "claude-opus-4-7",
        "gemini": "gemini-3-flash-preview",
        "google": "gemini-3-flash-preview",
    },
}

# Backwards-compat alias used by audio chain.
DEFAULT_MODELS = DEFAULT_MODELS_BY_CAPABILITY["audio"]

# Priority chain per capability. Lower number = tried first. Mirrors
# OpenClaw's `autoPriority` field per (provider, capability).
AUTO_PRIORITY_BY_CAPABILITY: dict[str, dict[str, int]] = {
    "audio": {
        "openai": 10,
        "groq": 20,
        "deepgram": 30,
        "gemini": 40,
        "google": 40,
        "mistral": 50,
        "xai": 60,
    },
    "image": {
        "openai": 10,
        "anthropic": 20,
        "gemini": 30,
        "google": 30,
        "minimax": 40,
        "qwen": 45,
        "moonshot": 50,
        "minimax-portal": 55,
        "zai": 60,
        "openrouter": 70,
        "xai": 80,
        # NOTE: openai-codex intentionally has NO autoPriority — matches
        # OpenClaw bundled-defaults.ts (codex is only used when user
        # explicitly sets active model to openai-codex/*). Active-chat-
        # model fallback (Tier 0) handles that case.
    },
    "video": {
        "gemini": 10,
        "google": 10,
        "qwen": 15,
        "moonshot": 20,
    },
    "document": {
        # Native PDF passthrough first (no extraction, model reads
        # PDF directly via inline_data); extraction is a separate
        # fallback path that runs INSIDE the hook without touching
        # this priority list.
        "anthropic": 10,
        "gemini": 20,
        "google": 20,
    },
}

# Backwards-compat alias for audio.
AUTO_PRIORITY = AUTO_PRIORITY_BY_CAPABILITY["audio"]

# Native HTTP endpoints. Per-capability where the provider exposes
# different API roots (e.g. MiniMax uses a special vlm endpoint for
# images). `OPENAI_COMPAT_BASE_URL` is what OpenAI-compatible
# providers use for both audio multipart AND chat.completions; the
# helpers below append the right path.
PROVIDER_BASE_URLS: dict[str, str] = {
    "openai": "https://api.openai.com/v1",
    "openai-codex": "https://api.openai.com/v1",  # codex uses OpenAI's API base
    "groq": "https://api.groq.com/openai/v1",
    "mistral": "https://api.mistral.ai/v1",
    "xai": "https://api.x.ai/v1",
    "deepgram": "https://api.deepgram.com/v1",
    "gemini": "https://generativelanguage.googleapis.com/v1beta",
    "google": "https://generativelanguage.googleapis.com/v1beta",
    "anthropic": "https://api.anthropic.com/v1",
    "minimax": "https://api.minimax.io",
    "minimax-portal": "https://api.minimax.io",  # same endpoint, different auth key
    "zai": "https://api.z.ai/api/paas/v4",
    "qwen": "https://dashscope.aliyuncs.com/compatible-mode/v1",
    "moonshot": "https://api.moonshot.cn/v1",
    "openrouter": "https://openrouter.ai/api/v1",
}

PROVIDER_KEY_ENV_VARS: dict[str, tuple[str, ...]] = {
    "openai": ("OPENAI_API_KEY",),
    # openai-codex uses Codex CLI OAuth in OpenClaw. For AgentBuff
    # mass-market we accept either dedicated key OR fall back to
    # OPENAI_API_KEY — chief's user won't have separate Codex
    # credentials.
    "openai-codex": ("OPENAI_CODEX_API_KEY", "OPENAI_API_KEY"),
    "groq": ("GROQ_API_KEY",),
    "deepgram": ("DEEPGRAM_API_KEY",),
    "gemini": ("GEMINI_API_KEY", "GOOGLE_API_KEY", "HERMES_DEFAULT_GEMINI_KEY"),
    "google": ("GOOGLE_API_KEY", "GEMINI_API_KEY", "HERMES_DEFAULT_GEMINI_KEY"),
    "mistral": ("MISTRAL_API_KEY",),
    "xai": ("XAI_API_KEY", "GROK_API_KEY"),
    "anthropic": ("ANTHROPIC_API_KEY",),
    "minimax": ("MINIMAX_API_KEY",),
    # minimax-portal is OpenClaw's secondary minimax channel (different
    # entitlements, different key in OpenClaw config). Try its dedicated
    # key first; fall back to MINIMAX_API_KEY.
    "minimax-portal": ("MINIMAX_PORTAL_API_KEY", "MINIMAX_API_KEY"),
    "zai": ("ZAI_API_KEY", "Z_AI_API_KEY"),
    "qwen": ("QWEN_API_KEY", "DASHSCOPE_API_KEY"),
    "moonshot": ("MOONSHOT_API_KEY", "KIMI_API_KEY"),
    "openrouter": ("OPENROUTER_API_KEY",),
}

_GEMINI_AUDIO_MIME_MAP = {
    "audio/mp3": "audio/mp3",
    "audio/mpeg": "audio/mp3",
    "audio/ogg": "audio/ogg",
    "audio/opus": "audio/ogg",
    "audio/oga": "audio/ogg",
    "audio/wav": "audio/wav",
    "audio/x-wav": "audio/wav",
    "audio/wave": "audio/wav",
    "audio/x-m4a": "audio/aac",
    "audio/mp4": "audio/aac",
    "audio/aac": "audio/aac",
    "audio/webm": "audio/webm",
    "audio/flac": "audio/flac",
}

_OPENAI_AUDIO_FILENAMES = {
    "audio/mp3": "audio.mp3",
    "audio/mpeg": "audio.mp3",
    "audio/ogg": "audio.ogg",
    "audio/opus": "audio.ogg",
    "audio/wav": "audio.wav",
    "audio/x-wav": "audio.wav",
    "audio/wave": "audio.wav",
    "audio/x-m4a": "audio.m4a",
    "audio/mp4": "audio.m4a",
    "audio/aac": "audio.m4a",
    "audio/webm": "audio.webm",
    "audio/flac": "audio.flac",
}

_ACTIVE_PROVIDER_ALIASES = {
    "google": "google",
    "gemini": "gemini",
    "openai": "openai",
    "openai-codex": "openai-codex",
    "codex": "openai-codex",
    "codex-cli": "openai-codex",
    "anthropic": "anthropic",
    "claude": "anthropic",
    "groq": "groq",
    "deepgram": "deepgram",
    "mistral": "mistral",
    "xai": "xai",
    "grok": "xai",
    "openrouter": "openrouter",
    "openrouter-direct": "openrouter",
    "kilocode": "openai",
    "ai-gateway": "openai",
    "minimax": "minimax",
    "minimax-portal": "minimax-portal",
    "zai": "zai",
    "z-ai": "zai",
    "glm": "zai",
    "qwen": "qwen",
    "dashscope": "qwen",
    "moonshot": "moonshot",
    "kimi": "moonshot",
}


# ──────────────────────────────────────────────────────────────────────
# Env + config resolution (no Hermes imports — Hermes-update-safe)
# ──────────────────────────────────────────────────────────────────────


def _load_dot_env_file() -> dict[str, str]:
    cached = getattr(_load_dot_env_file, "_cache", None)
    if cached is not None:
        return cached  # type: ignore[return-value]
    out: dict[str, str] = {}
    home = os.environ.get("HERMES_HOME") or str(Path.home() / ".hermes")
    env_path = Path(home) / ".env"
    try:
        if env_path.is_file():
            for raw in env_path.read_text(
                encoding="utf-8", errors="replace"
            ).splitlines():
                line = raw.strip()
                if not line or line.startswith("#") or "=" not in line:
                    continue
                key, _, value = line.partition("=")
                key = key.strip()
                if key.startswith("export "):
                    key = key[len("export "):].strip()
                if not key:
                    continue
                value = value.strip()
                if (
                    len(value) >= 2
                    and value[0] == value[-1]
                    and value[0] in ("'", '"')
                ):
                    value = value[1:-1]
                out[key] = value
    except Exception as exc:
        logger.debug("agentbuff-multimodal: read .env failed: %s", exc)
    setattr(_load_dot_env_file, "_cache", out)
    return out


def _resolve_env_value(*var_names: str) -> Optional[str]:
    for name in var_names:
        v = os.environ.get(name)
        if v and v.strip():
            return v.strip()
    dot_env = _load_dot_env_file()
    for name in var_names:
        v = dot_env.get(name)
        if v and v.strip():
            return v.strip()
    return None


def _get_provider_api_key(provider_id: str) -> Optional[str]:
    var_names = PROVIDER_KEY_ENV_VARS.get(provider_id.lower())
    if not var_names:
        return None
    return _resolve_env_value(*var_names)


def _get_active_chat_provider() -> Optional[str]:
    """Read `~/.hermes/config.yaml::model.default` and extract provider prefix."""
    home = os.environ.get("HERMES_HOME") or str(Path.home() / ".hermes")
    cfg_path = Path(home) / "config.yaml"
    if not cfg_path.is_file():
        return None
    try:
        text = cfg_path.read_text(encoding="utf-8", errors="replace")
    except Exception:
        return None
    match = re.search(r"^model\s*:\s*\n((?:\s+.+\n?)+)", text, flags=re.MULTILINE)
    if match:
        block = match.group(1)
        sub = re.search(
            r"^\s+(?:default|model)\s*:\s*['\"]?([^'\"\n#]+)['\"]?",
            block,
            flags=re.MULTILINE,
        )
        candidate = sub.group(1).strip() if sub else None
    else:
        inline = re.search(
            r"^model\s*:\s*['\"]?([^'\"\n#]+)['\"]?", text, flags=re.MULTILINE
        )
        candidate = inline.group(1).strip() if inline else None
    if not candidate or "/" not in candidate:
        return None
    prefix = candidate.split("/", 1)[0].strip().lower()
    return _ACTIVE_PROVIDER_ALIASES.get(prefix)


# ══════════════════════════════════════════════════════════════════════
# Capability-disable config — `tools.media.<cap>.enabled` flag
# ══════════════════════════════════════════════════════════════════════
#
# Mirrors OpenClaw's `cfg.tools?.media?.<cap>?.enabled` check
# (`media-understanding/apply.ts:6420+`). When set to `false`, the
# corresponding chain is short-circuited — useful for admin control
# (e.g. disable expensive video on cheap subscription tier).
#
# Config shape in `~/.hermes/config.yaml`:
#   tools:
#     media:
#       audio:
#         enabled: false   # disable STT entirely
#         echoTranscript: true   # send 📝 "..." preview msg
#       image:
#         enabled: true
#       video:
#         enabled: true
#       document:
#         enabled: true
#
# Default = enabled. Missing config means feature is ON.


def _read_yaml_tools_media_value(*keys: str) -> Optional[Any]:
    """Read a value from `~/.hermes/config.yaml::tools.media.<keys>`.

    Pure-Python minimal YAML walker (avoids pyyaml dep at import-time;
    we only need leaf scalars under a known nested path). Returns
    `None` if any key in the chain is missing.
    """
    home = os.environ.get("HERMES_HOME") or str(Path.home() / ".hermes")
    cfg_path = Path(home) / "config.yaml"
    if not cfg_path.is_file():
        return None
    try:
        import yaml  # PyYAML is in bridge requirements + hermes-agent deps
        with open(cfg_path, "r", encoding="utf-8") as f:
            data = yaml.safe_load(f) or {}
    except Exception:
        return None
    if not isinstance(data, dict):
        return None
    node: Any = data.get("tools")
    if not isinstance(node, dict):
        return None
    node = node.get("media")
    if not isinstance(node, dict):
        return None
    for key in keys:
        if not isinstance(node, dict):
            return None
        node = node.get(key)
        if node is None:
            return None
    return node


def _capability_enabled(capability: str) -> bool:
    """True unless `tools.media.<capability>.enabled` is explicitly false."""
    value = _read_yaml_tools_media_value(capability, "enabled")
    if value is False:
        return False
    return True


def _capability_echo_transcript() -> bool:
    """True if `tools.media.audio.echoTranscript: true`. Default false."""
    value = _read_yaml_tools_media_value("audio", "echoTranscript")
    return value is True


def _capability_audio_preflight() -> bool:
    """True if `tools.media.audio.preflight: true`. Default false.
    When true, hook transcribes audio in pre_gateway_dispatch so the
    transcript appears in event.text before any auth/mention check."""
    value = _read_yaml_tools_media_value("audio", "preflight")
    return value is True


# ══════════════════════════════════════════════════════════════════════
# Binary cache — avoid re-reading same file when both hook AND Hermes
# enrich call transcribe_audio for the same path.
# ══════════════════════════════════════════════════════════════════════
#
# Keyed on (path, mtime, size). Each entry caches:
#   * raw bytes (saves disk read on retry)
#   * transcribe_audio result dict (saves provider API call when both
#     hook + Hermes enrich pass the same path through the chain)
#
# LRU semantics via collections.OrderedDict — capped at 32 entries
# to bound memory. Audio files cap at MAX_AUDIO_BYTES = 10MB so worst
# case ~320MB which is fine for a container.


class _BinaryCache:
    """LRU file-bytes + transcribe-result cache."""

    def __init__(self, max_entries: int = 32) -> None:
        import collections
        self._max = max_entries
        self._bytes: "collections.OrderedDict[tuple, bytes]" = collections.OrderedDict()
        self._transcript: "collections.OrderedDict[tuple, dict]" = collections.OrderedDict()

    def _key(self, file_path: str) -> Optional[tuple]:
        try:
            st = Path(file_path).stat()
            return (str(Path(file_path).resolve()), st.st_mtime_ns, st.st_size)
        except Exception:
            return None

    def get_bytes(self, file_path: str) -> Optional[bytes]:
        key = self._key(file_path)
        if key is None or key not in self._bytes:
            return None
        self._bytes.move_to_end(key)
        return self._bytes[key]

    def put_bytes(self, file_path: str, data: bytes) -> None:
        key = self._key(file_path)
        if key is None:
            return
        self._bytes[key] = data
        self._bytes.move_to_end(key)
        while len(self._bytes) > self._max:
            self._bytes.popitem(last=False)

    def get_transcript(self, file_path: str) -> Optional[dict]:
        key = self._key(file_path)
        if key is None or key not in self._transcript:
            return None
        self._transcript.move_to_end(key)
        return self._transcript[key]

    def put_transcript(self, file_path: str, result: dict) -> None:
        key = self._key(file_path)
        if key is None:
            return
        self._transcript[key] = result
        self._transcript.move_to_end(key)
        while len(self._transcript) > self._max:
            self._transcript.popitem(last=False)


_BINARY_CACHE = _BinaryCache()


def _read_file_cached(file_path: str) -> Optional[bytes]:
    """Read file bytes via cache. Returns None on read error."""
    cached = _BINARY_CACHE.get_bytes(file_path)
    if cached is not None:
        return cached
    try:
        data = Path(file_path).read_bytes()
    except Exception:
        return None
    _BINARY_CACHE.put_bytes(file_path, data)
    return data


# ══════════════════════════════════════════════════════════════════════
# Concurrent execution — run multiple media chains in parallel
# ══════════════════════════════════════════════════════════════════════
#
# Mirrors OpenClaw's `media-understanding/concurrency.ts`. When the
# hook receives multiple attachments in one event, we fan out the
# chain calls via ThreadPoolExecutor so total wall-clock matches the
# SLOWEST provider rather than the SUM.
#
# Limit configurable via env `MEDIA_CONCURRENCY` (default 3).
# httpx.Client is thread-safe for read; each thread does its own POST.


def _media_concurrency_limit() -> int:
    raw = os.environ.get("MEDIA_CONCURRENCY") or os.environ.get("HERMES_MEDIA_CONCURRENCY")
    if not raw:
        return 3
    try:
        n = int(raw)
        return max(1, min(n, 16))
    except (TypeError, ValueError):
        return 3


def _run_concurrently(tasks: list[Callable[[], Any]]) -> list[Any]:
    """Run callables in a thread pool, return results in original order.
    Exceptions are caught and replaced with the exception object."""
    if not tasks:
        return []
    if len(tasks) == 1:
        try:
            return [tasks[0]()]
        except Exception as exc:
            return [exc]
    from concurrent.futures import ThreadPoolExecutor
    limit = min(len(tasks), _media_concurrency_limit())
    with ThreadPoolExecutor(max_workers=limit) as pool:
        futures = [pool.submit(task) for task in tasks]
        return [
            (f.result() if not f.exception() else f.exception())
            for f in futures
        ]


# ══════════════════════════════════════════════════════════════════════
# CLI fallback — spawn local binary (e.g. whisper.cpp)
# ══════════════════════════════════════════════════════════════════════
#
# Mirrors OpenClaw's `type: "cli"` decision type. When the user has a
# local whisper binary installed (whisper.cpp, OpenAI-whisper Python
# CLI, faster-whisper CLI), we shell out to it as Tier 95 — between
# our HTTP provider chain (Tier 0+1) and Hermes' built-in `transcribe_
# audio` (Tier 99).
#
# Config: env var `MEDIA_CLI_WHISPER_CMD` = full command template.
# `{path}` placeholder = audio file path.
# Example: `MEDIA_CLI_WHISPER_CMD="/usr/bin/whisper {path} --model base --output_format txt --output_dir /tmp"`
#
# Plugin parses stdout for transcript. If the binary writes to a file
# instead of stdout, set `MEDIA_CLI_WHISPER_OUTPUT_PATTERN={path}.txt`
# and we'll read from that path.


def _transcribe_via_cli(file_path: str) -> tuple[Optional[str], Optional[str]]:
    cmd_template = os.environ.get("MEDIA_CLI_WHISPER_CMD")
    if not cmd_template:
        return (None, "no MEDIA_CLI_WHISPER_CMD configured")
    import shlex
    import subprocess
    try:
        cmd_str = cmd_template.replace("{path}", shlex.quote(file_path))
        cmd = shlex.split(cmd_str)
        timeout_s = int(os.environ.get("MEDIA_CLI_WHISPER_TIMEOUT", "120"))
        result = subprocess.run(
            cmd,
            capture_output=True,
            text=True,
            timeout=timeout_s,
            check=False,
        )
        if result.returncode != 0:
            err = (result.stderr or "")[:200]
            return (None, f"cli exit {result.returncode}: {err}")
        # Strategy 1: read from a sidecar file if pattern configured
        out_pattern = os.environ.get("MEDIA_CLI_WHISPER_OUTPUT_PATTERN")
        if out_pattern:
            out_path = out_pattern.replace("{path}", file_path)
            try:
                text = Path(out_path).read_text(encoding="utf-8", errors="replace")
                cleaned = text.strip()
                if cleaned:
                    return (cleaned, None)
            except Exception:
                pass
        # Strategy 2: parse stdout
        stdout = (result.stdout or "").strip()
        if stdout:
            return (stdout, None)
        return (None, "cli produced no output")
    except subprocess.TimeoutExpired:
        return (None, "cli timeout")
    except Exception as exc:
        return (None, f"{type(exc).__name__}: {exc}")


# ══════════════════════════════════════════════════════════════════════
# Provider implementations
# ══════════════════════════════════════════════════════════════════════
#
# Each provider's HTTP call goes through `_http_post()` which consults
# config.yaml for per-provider overrides (TLS, proxy, baseUrl, auth
# header). Mirrors OpenClaw's `MediaUnderstandingProviderRequestAuth/
# Tls/ProxyOverride` types from `media-understanding/types.ts:63-79`.
#
# Override config shape in `~/.hermes/config.yaml`:
#   tools:
#     media:
#       providers:
#         openai:
#           baseUrl: "https://my-custom-openai.example.com/v1"
#           proxy: "http://proxy.corp:8080"
#           tlsVerify: false           # skip cert validation (testing only)
#           tlsCaCert: "/etc/ssl/corporate-ca.pem"
#           authHeader: { name: "x-api-key", value: "..." }


def _read_provider_overrides(provider_id: str) -> dict:
    """Read per-provider overrides from tools.media.providers.<id>.* in config.yaml.
    Returns empty dict if missing."""
    overrides = _read_yaml_tools_media_value("providers", provider_id)
    return overrides if isinstance(overrides, dict) else {}


def _http_post(url: str, *, provider_id: Optional[str] = None, **kwargs):
    """HTTP POST with optional per-provider TLS/proxy/authHeader override.

    Lazy-imports httpx so the plugin can load even if httpx is briefly
    unavailable in a Hermes container build.
    """
    import httpx
    timeout = httpx.Timeout(connect=10.0, read=90.0, write=90.0, pool=10.0)

    client_kwargs: dict[str, Any] = {"timeout": timeout}
    if provider_id:
        overrides = _read_provider_overrides(provider_id)
        if overrides:
            # TLS overrides
            tls_verify = overrides.get("tlsVerify")
            if tls_verify is False:
                client_kwargs["verify"] = False
            elif tls_verify is True:
                client_kwargs["verify"] = True
            tls_ca = overrides.get("tlsCaCert")
            if isinstance(tls_ca, str) and tls_ca.strip():
                client_kwargs["verify"] = tls_ca.strip()
            # Proxy override
            proxy_url = overrides.get("proxy")
            if isinstance(proxy_url, str) and proxy_url.strip():
                # httpx 0.28 uses `proxy=` (singular). Older `proxies=` also accepted.
                client_kwargs["proxy"] = proxy_url.strip()
            # Custom auth header (overrides whatever was passed by caller)
            auth = overrides.get("authHeader")
            if isinstance(auth, dict):
                name = auth.get("name")
                value = auth.get("value")
                if isinstance(name, str) and isinstance(value, str):
                    headers = kwargs.setdefault("headers", {})
                    if isinstance(headers, dict):
                        headers[name] = value

    with httpx.Client(**client_kwargs) as client:
        return client.post(url, **kwargs)


def _transcribe_openai_compatible(
    provider_id: str, audio_bytes: bytes, mime: str
) -> tuple[Optional[str], Optional[str]]:
    api_key = _get_provider_api_key(provider_id)
    if not api_key:
        return (None, "no_api_key")
    base_url = (
        os.environ.get(f"{provider_id.upper()}_BASE_URL")
        or PROVIDER_BASE_URLS.get(provider_id)
    )
    if not base_url:
        return (None, "no base_url")
    model = DEFAULT_MODELS_BY_CAPABILITY["audio"].get(provider_id, "whisper-1")
    canon_mime = (mime or "").lower() or "audio/ogg"
    filename = _OPENAI_AUDIO_FILENAMES.get(canon_mime, "audio.ogg")
    try:
        files = {"file": (filename, io.BytesIO(audio_bytes), canon_mime)}
        data = {"model": model}
        resp = _http_post(
            f"{base_url.rstrip('/')}/audio/transcriptions",
            provider_id=provider_id,
            headers={"Authorization": f"Bearer {api_key}"},
            files=files,
            data=data,
        )
        if resp.status_code != 200:
            return (None, f"HTTP {resp.status_code}: {resp.text[:200]}")
        payload = resp.json()
        return ((payload.get("text") or "").strip(), None)
    except Exception as exc:
        return (None, f"{type(exc).__name__}: {exc}")


def _transcribe_gemini(
    audio_bytes: bytes, mime: str
) -> tuple[Optional[str], Optional[str]]:
    api_key = _get_provider_api_key("gemini")
    if not api_key:
        return (None, "no_api_key")
    canon_mime = _GEMINI_AUDIO_MIME_MAP.get((mime or "").lower(), "audio/ogg")
    model = (
        os.environ.get("GEMINI_STT_MODEL")
        or DEFAULT_MODELS_BY_CAPABILITY["audio"]["gemini"]
    )
    try:
        b64 = base64.b64encode(audio_bytes).decode("ascii")
        body = {
            "contents": [
                {
                    "parts": [
                        {
                            "text": (
                                "Transcribe this audio exactly as spoken. "
                                "Output ONLY the verbatim transcript — no commentary, "
                                "no quotation marks, no language detection prefix, "
                                "no timestamps, no speaker labels. Preserve the "
                                "speaker's original language (do NOT translate)."
                            )
                        },
                        {
                            "inline_data": {
                                "mime_type": canon_mime,
                                "data": b64,
                            }
                        },
                    ]
                }
            ],
            "generationConfig": {"temperature": 0.0, "maxOutputTokens": 4096},
        }
        url = (
            f"{PROVIDER_BASE_URLS['gemini']}/models/{model}:generateContent?key={api_key}"
        )
        resp = _http_post(url, provider_id="gemini", json=body)
        if resp.status_code != 200:
            return (None, f"HTTP {resp.status_code}: {resp.text[:200]}")
        data = resp.json()
        candidates = data.get("candidates") or []
        if not candidates:
            return (None, "no candidates")
        parts = (candidates[0].get("content") or {}).get("parts") or []
        text = "".join(
            (p.get("text") or "") for p in parts if isinstance(p, dict)
        ).strip()
        return (text, None)
    except Exception as exc:
        return (None, f"{type(exc).__name__}: {exc}")


def _transcribe_deepgram(
    audio_bytes: bytes, mime: str
) -> tuple[Optional[str], Optional[str]]:
    api_key = _get_provider_api_key("deepgram")
    if not api_key:
        return (None, "no_api_key")
    model = DEFAULT_MODELS_BY_CAPABILITY["audio"]["deepgram"]
    canon_mime = (mime or "").lower() or "audio/ogg"
    url = (
        f"{PROVIDER_BASE_URLS['deepgram']}/listen"
        f"?model={model}&smart_format=true&punctuate=true"
    )
    try:
        resp = _http_post(
            url,
            provider_id="deepgram",
            headers={
                "Authorization": f"Token {api_key}",
                "Content-Type": canon_mime,
            },
            content=audio_bytes,
        )
        if resp.status_code != 200:
            return (None, f"HTTP {resp.status_code}: {resp.text[:200]}")
        data = resp.json()
        channels = (data.get("results") or {}).get("channels") or []
        if not channels:
            return ("", None)
        alternatives = channels[0].get("alternatives") or []
        if not alternatives:
            return ("", None)
        return ((alternatives[0].get("transcript") or "").strip(), None)
    except Exception as exc:
        return (None, f"{type(exc).__name__}: {exc}")


PROVIDER_TRANSCRIBE_FNS: dict[
    str, Callable[[bytes, str], tuple[Optional[str], Optional[str]]]
] = {
    "openai": lambda b, m: _transcribe_openai_compatible("openai", b, m),
    "groq": lambda b, m: _transcribe_openai_compatible("groq", b, m),
    "mistral": lambda b, m: _transcribe_openai_compatible("mistral", b, m),
    "xai": lambda b, m: _transcribe_openai_compatible("xai", b, m),
    "gemini": _transcribe_gemini,
    "google": _transcribe_gemini,
    "deepgram": _transcribe_deepgram,
}


def _detect_mime_from_path(file_path: str) -> str:
    """Best-effort MIME from extension. Hermes' channel adapters set
    `event.media_types` separately, but `transcribe_audio` only gets a
    path, so we re-derive."""
    ext = Path(file_path).suffix.lower()
    mapping = {
        ".mp3": "audio/mp3", ".mpeg": "audio/mp3", ".mpga": "audio/mp3",
        ".m4a": "audio/x-m4a", ".mp4": "audio/mp4",
        ".aac": "audio/aac",
        ".ogg": "audio/ogg", ".oga": "audio/ogg", ".opus": "audio/opus",
        ".wav": "audio/wav",
        ".webm": "audio/webm",
        ".flac": "audio/flac",
    }
    return mapping.get(ext, "audio/ogg")


def _run_extended_chain(
    file_path: str,
) -> tuple[Optional[str], Optional[str], list[str]]:
    """Read the audio file and run our extended provider chain.

    Returns (transcript, provider_used, attempts). transcript=None means
    no provider produced a usable result. attempts is the diagnostic log.
    """
    # Capability disable check (Tier -1)
    if not _capability_enabled("audio"):
        return (None, None, ["audio capability disabled in config"])

    audio_bytes = _read_file_cached(file_path)
    if audio_bytes is None:
        return (None, None, [f"read_file: failed to read {file_path}"])
    mime = _detect_mime_from_path(file_path)
    attempts: list[str] = []
    tried: set[str] = set()

    def _canon(pid: str) -> str:
        return "gemini" if pid == "google" else pid

    def _try(provider_id: str) -> tuple[Optional[str], Optional[str]]:
        canonical = _canon(provider_id)
        if canonical in tried:
            return (None, "already_tried")
        tried.add(canonical)
        fn = PROVIDER_TRANSCRIBE_FNS.get(provider_id)
        if not fn:
            attempts.append(f"{provider_id}: not registered")
            return (None, "not_registered")
        text, err = fn(audio_bytes, mime)
        if text is not None:
            attempts.append(f"{provider_id}: ok ({len(text)} chars)")
        else:
            attempts.append(f"{provider_id}: {err}")
        return (text, err)

    # Tier 0: active chat-model provider
    active = _get_active_chat_provider()
    if active:
        text, _ = _try(active)
        if text is not None:
            return (text, active, attempts)

    # Tier 1+: bundled priority chain
    items = sorted(
        {k: AUTO_PRIORITY[k] for k in AUTO_PRIORITY}.items(),
        key=lambda kv: (kv[1], kv[0]),
    )
    for provider_id, _prio in items:
        text, _ = _try(provider_id)
        if text is not None:
            return (text, provider_id, attempts)

    # Tier 95: CLI fallback (local whisper.cpp / OpenAI-whisper CLI)
    # Runs only if `MEDIA_CLI_WHISPER_CMD` env is set. Useful for
    # offline-first deployments where chief doesn't want API calls.
    cli_text, cli_err = _transcribe_via_cli(file_path)
    if cli_text:
        attempts.append(f"cli: ok ({len(cli_text)} chars)")
        return (cli_text, "cli", attempts)
    if cli_err and "no MEDIA_CLI_WHISPER_CMD" not in cli_err:
        attempts.append(f"cli: {cli_err}")

    return (None, None, attempts)


# ══════════════════════════════════════════════════════════════════════
# VISION (IMAGE) chain — describe images via 9 providers + active model
# ══════════════════════════════════════════════════════════════════════
#
# Strategy: each provider has its own HTTP wire shape. We re-implement
# OpenClaw's `media-understanding/image.ts` + provider-specific code
# (notably `agents/minimax-vlm.ts`) in Python. Common pattern:
#
#   describe_<provider>(image_bytes, mime, prompt) → (text_or_None, err)
#
# Three families:
#   * OpenAI-compatible chat/completions with image_url data-URL part:
#     openai, openrouter, xai, qwen, moonshot, zai (all expose the
#     same "vision in chat.completions" shape).
#   * Anthropic native messages API with image source { type:"base64" }.
#   * Gemini native generateContent with inline_data part.
#   * MiniMax custom /v1/coding_plan/vlm endpoint.
#
# All return PLAIN TEXT description (no JSON envelope). The wrapper
# around `vision_analyze_tool` turns it into Hermes' contract shape.


_IMAGE_MIME_MAP = {
    "image/jpg": "image/jpeg",
    "image/jpeg": "image/jpeg",
    "image/png": "image/png",
    "image/gif": "image/gif",
    "image/webp": "image/webp",
    "image/bmp": "image/bmp",
}


def _detect_image_mime_from_path(file_path: str) -> str:
    ext = Path(file_path).suffix.lower()
    mapping = {
        ".jpg": "image/jpeg",
        ".jpeg": "image/jpeg",
        ".png": "image/png",
        ".gif": "image/gif",
        ".webp": "image/webp",
        ".bmp": "image/bmp",
    }
    return mapping.get(ext, "image/jpeg")


def _vision_default_prompt(user_prompt: Optional[str]) -> str:
    if user_prompt and user_prompt.strip():
        return user_prompt.strip()
    return (
        "Describe this image in detail. Note any text visible in the image "
        "verbatim. If it's a document or screenshot, transcribe the content. "
        "If it's a chart, describe data trends. Be thorough but concise."
    )


def _describe_openai_compatible_image(
    provider_id: str,
    image_bytes: bytes,
    mime: str,
    prompt: str,
) -> tuple[Optional[str], Optional[str]]:
    """OpenAI chat.completions shape with image_url data-URL parts.

    Covers: openai, openrouter, xai, qwen (Dashscope compatible-mode),
    moonshot, zai.
    """
    api_key = _get_provider_api_key(provider_id)
    if not api_key:
        return (None, "no_api_key")
    base_url = (
        os.environ.get(f"{provider_id.upper()}_BASE_URL")
        or PROVIDER_BASE_URLS.get(provider_id)
    )
    if not base_url:
        return (None, "no base_url")
    model = (
        os.environ.get(f"{provider_id.upper()}_VISION_MODEL")
        or DEFAULT_MODELS_BY_CAPABILITY["image"].get(provider_id)
    )
    if not model:
        return (None, "no default vision model")
    canon_mime = _IMAGE_MIME_MAP.get((mime or "").lower(), "image/jpeg")
    try:
        b64 = base64.b64encode(image_bytes).decode("ascii")
        body = {
            "model": model,
            "messages": [
                {
                    "role": "user",
                    "content": [
                        {"type": "text", "text": prompt},
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:{canon_mime};base64,{b64}",
                            },
                        },
                    ],
                }
            ],
            "max_tokens": 1024,
        }
        resp = _http_post(
            f"{base_url.rstrip('/')}/chat/completions",
            provider_id=provider_id,
            headers={
                "Authorization": f"Bearer {api_key}",
                "Content-Type": "application/json",
            },
            json=body,
        )
        if resp.status_code != 200:
            return (None, f"HTTP {resp.status_code}: {resp.text[:200]}")
        data = resp.json()
        choices = data.get("choices") or []
        if not choices:
            return (None, "no choices")
        message = choices[0].get("message") or {}
        content = message.get("content")
        if isinstance(content, str):
            return (content.strip(), None)
        if isinstance(content, list):
            text = "\n".join(
                (p.get("text") or "").strip()
                for p in content
                if isinstance(p, dict) and p.get("text")
            ).strip()
            if text:
                return (text, None)
        reasoning = message.get("reasoning_content")
        if isinstance(reasoning, str) and reasoning.strip():
            return (reasoning.strip(), None)
        return (None, "empty content")
    except Exception as exc:
        return (None, f"{type(exc).__name__}: {exc}")


def _describe_anthropic_image(
    image_bytes: bytes, mime: str, prompt: str
) -> tuple[Optional[str], Optional[str]]:
    """Anthropic Messages API native vision (claude vision)."""
    api_key = _get_provider_api_key("anthropic")
    if not api_key:
        return (None, "no_api_key")
    model = (
        os.environ.get("ANTHROPIC_VISION_MODEL")
        or DEFAULT_MODELS_BY_CAPABILITY["image"]["anthropic"]
    )
    canon_mime = _IMAGE_MIME_MAP.get((mime or "").lower(), "image/jpeg")
    try:
        b64 = base64.b64encode(image_bytes).decode("ascii")
        body = {
            "model": model,
            "max_tokens": 1024,
            "messages": [
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "image",
                            "source": {
                                "type": "base64",
                                "media_type": canon_mime,
                                "data": b64,
                            },
                        },
                        {"type": "text", "text": prompt},
                    ],
                }
            ],
        }
        resp = _http_post(
            f"{PROVIDER_BASE_URLS['anthropic']}/messages",
            provider_id="anthropic",
            headers={
                "x-api-key": api_key,
                "anthropic-version": "2023-06-01",
                "content-type": "application/json",
            },
            json=body,
        )
        if resp.status_code != 200:
            return (None, f"HTTP {resp.status_code}: {resp.text[:200]}")
        data = resp.json()
        content = data.get("content") or []
        text = "\n".join(
            (p.get("text") or "").strip()
            for p in content
            if isinstance(p, dict) and p.get("type") == "text"
        ).strip()
        return ((text or None), None if text else "empty content")
    except Exception as exc:
        return (None, f"{type(exc).__name__}: {exc}")


def _describe_gemini_image(
    image_bytes: bytes, mime: str, prompt: str
) -> tuple[Optional[str], Optional[str]]:
    """Gemini generateContent with inline_data part."""
    api_key = _get_provider_api_key("gemini")
    if not api_key:
        return (None, "no_api_key")
    model = (
        os.environ.get("GEMINI_VISION_MODEL")
        or DEFAULT_MODELS_BY_CAPABILITY["image"]["gemini"]
    )
    canon_mime = _IMAGE_MIME_MAP.get((mime or "").lower(), "image/jpeg")
    try:
        b64 = base64.b64encode(image_bytes).decode("ascii")
        body = {
            "contents": [
                {
                    "parts": [
                        {"text": prompt},
                        {
                            "inline_data": {
                                "mime_type": canon_mime,
                                "data": b64,
                            }
                        },
                    ]
                }
            ],
            "generationConfig": {"temperature": 0.1, "maxOutputTokens": 1024},
        }
        url = (
            f"{PROVIDER_BASE_URLS['gemini']}/models/{model}:generateContent"
            f"?key={api_key}"
        )
        resp = _http_post(url, provider_id="gemini", json=body)
        if resp.status_code != 200:
            return (None, f"HTTP {resp.status_code}: {resp.text[:200]}")
        data = resp.json()
        candidates = data.get("candidates") or []
        if not candidates:
            return (None, "no candidates")
        parts = (candidates[0].get("content") or {}).get("parts") or []
        text = "".join(
            (p.get("text") or "") for p in parts if isinstance(p, dict)
        ).strip()
        return ((text or None), None if text else "empty content")
    except Exception as exc:
        return (None, f"{type(exc).__name__}: {exc}")


def _describe_minimax_image(
    image_bytes: bytes, mime: str, prompt: str
) -> tuple[Optional[str], Optional[str]]:
    """MiniMax custom VLM endpoint (matches OpenClaw's agents/minimax-vlm.ts).

    Wire shape:
        POST {host}/v1/coding_plan/vlm
        Authorization: Bearer <key>
        MM-API-Source: AgentBuff
        body: {prompt, image_url: "data:image/...;base64,..."}
    """
    api_key = _get_provider_api_key("minimax")
    if not api_key:
        return (None, "no_api_key")
    canon_mime = _IMAGE_MIME_MAP.get((mime or "").lower(), "image/jpeg")
    base_url = (
        os.environ.get("MINIMAX_BASE_URL")
        or PROVIDER_BASE_URLS["minimax"]
    )
    try:
        b64 = base64.b64encode(image_bytes).decode("ascii")
        body = {
            "prompt": prompt,
            "image_url": f"data:{canon_mime};base64,{b64}",
        }
        resp = _http_post(
            f"{base_url.rstrip('/')}/v1/coding_plan/vlm",
            provider_id="minimax",
            headers={
                "Authorization": f"Bearer {api_key}",
                "Content-Type": "application/json",
                "MM-API-Source": "AgentBuff",
            },
            json=body,
        )
        if resp.status_code != 200:
            return (None, f"HTTP {resp.status_code}: {resp.text[:200]}")
        data = resp.json()
        base_resp = data.get("base_resp") or {}
        status_code = base_resp.get("status_code")
        if status_code not in (None, 0):
            return (
                None,
                f"minimax error {status_code}: {base_resp.get('status_msg')}",
            )
        content = data.get("content")
        if isinstance(content, str) and content.strip():
            return (content.strip(), None)
        return (None, "empty content")
    except Exception as exc:
        return (None, f"{type(exc).__name__}: {exc}")


PROVIDER_DESCRIBE_IMAGE_FNS: dict[
    str,
    Callable[[bytes, str, str], tuple[Optional[str], Optional[str]]],
] = {
    "openai": lambda b, m, p: _describe_openai_compatible_image("openai", b, m, p),
    "openai-codex": lambda b, m, p: _describe_openai_compatible_image(
        "openai-codex", b, m, p
    ),
    "openrouter": lambda b, m, p: _describe_openai_compatible_image(
        "openrouter", b, m, p
    ),
    "xai": lambda b, m, p: _describe_openai_compatible_image("xai", b, m, p),
    "qwen": lambda b, m, p: _describe_openai_compatible_image("qwen", b, m, p),
    "moonshot": lambda b, m, p: _describe_openai_compatible_image(
        "moonshot", b, m, p
    ),
    "zai": lambda b, m, p: _describe_openai_compatible_image("zai", b, m, p),
    "anthropic": _describe_anthropic_image,
    "gemini": _describe_gemini_image,
    "google": _describe_gemini_image,
    "minimax": _describe_minimax_image,
    # minimax-portal uses the same VLM endpoint shape as minimax but
    # with a different API key namespace. We pass provider_id through
    # so _get_provider_api_key picks the right env var.
    "minimax-portal": lambda b, m, p: _describe_minimax_image_with_provider(
        "minimax-portal", b, m, p
    ),
}


def _describe_minimax_image_with_provider(
    provider_id: str, image_bytes: bytes, mime: str, prompt: str
) -> tuple[Optional[str], Optional[str]]:
    """Same as `_describe_minimax_image` but parameterised on provider_id
    so we can distinguish `minimax` vs `minimax-portal` (different API
    keys, possibly different endpoints in future)."""
    api_key = _get_provider_api_key(provider_id)
    if not api_key:
        return (None, "no_api_key")
    canon_mime = _IMAGE_MIME_MAP.get((mime or "").lower(), "image/jpeg")
    base_url = (
        os.environ.get(f"{provider_id.upper().replace('-', '_')}_BASE_URL")
        or PROVIDER_BASE_URLS.get(provider_id)
    )
    if not base_url:
        return (None, "no base_url")
    try:
        b64 = base64.b64encode(image_bytes).decode("ascii")
        body = {
            "prompt": prompt,
            "image_url": f"data:{canon_mime};base64,{b64}",
        }
        resp = _http_post(
            f"{base_url.rstrip('/')}/v1/coding_plan/vlm",
            provider_id=provider_id,
            headers={
                "Authorization": f"Bearer {api_key}",
                "Content-Type": "application/json",
                "MM-API-Source": "AgentBuff",
            },
            json=body,
        )
        if resp.status_code != 200:
            return (None, f"HTTP {resp.status_code}: {resp.text[:200]}")
        data = resp.json()
        base_resp = data.get("base_resp") or {}
        status_code = base_resp.get("status_code")
        if status_code not in (None, 0):
            return (
                None,
                f"{provider_id} error {status_code}: {base_resp.get('status_msg')}",
            )
        content = data.get("content")
        if isinstance(content, str) and content.strip():
            return (content.strip(), None)
        return (None, "empty content")
    except Exception as exc:
        return (None, f"{type(exc).__name__}: {exc}")


def _run_vision_chain(
    file_path: str, prompt: str
) -> tuple[Optional[str], Optional[str], list[str]]:
    """Read image file and run extended provider chain.

    Returns (description, provider_used, attempts). description=None →
    no provider produced output. Mirrors `_run_extended_chain` for
    audio.
    """
    if not _capability_enabled("image"):
        return (None, None, ["image capability disabled in config"])

    image_bytes = _read_file_cached(file_path)
    if image_bytes is None:
        return (None, None, [f"read_file: failed to read {file_path}"])
    mime = _detect_image_mime_from_path(file_path)
    attempts: list[str] = []
    tried: set[str] = set()

    def _canon(pid: str) -> str:
        return "gemini" if pid == "google" else pid

    def _try(provider_id: str) -> tuple[Optional[str], Optional[str]]:
        canonical = _canon(provider_id)
        if canonical in tried:
            return (None, "already_tried")
        tried.add(canonical)
        fn = PROVIDER_DESCRIBE_IMAGE_FNS.get(provider_id)
        if not fn:
            attempts.append(f"{provider_id}: not registered")
            return (None, "not_registered")
        text, err = fn(image_bytes, mime, prompt)
        if text:
            attempts.append(f"{provider_id}: ok ({len(text)} chars)")
        else:
            attempts.append(f"{provider_id}: {err or 'empty'}")
        return (text, err)

    # Tier 0: active chat-model provider
    active = _get_active_chat_provider()
    if active:
        text, _ = _try(active)
        if text:
            return (text, active, attempts)

    # Tier 1+: bundled priority chain
    priority = AUTO_PRIORITY_BY_CAPABILITY["image"]
    items = sorted(priority.items(), key=lambda kv: (kv[1], kv[0]))
    for provider_id, _prio in items:
        text, _ = _try(provider_id)
        if text:
            return (text, provider_id, attempts)

    return (None, None, attempts)


# ══════════════════════════════════════════════════════════════════════
# VIDEO chain — describe videos via 3 providers (Gemini/Qwen/Moonshot)
# ══════════════════════════════════════════════════════════════════════
#
# Uses OpenAI-compatible chat.completions with `type: "video_url"`
# parts (data-URL base64). Note: video files are LARGER than images
# (per-bridge MAX_VIDEO_BYTES = 25MB, but providers often cap at lower
# inline limits). Gemini also supports the File API for big videos,
# but we use inline_data for simplicity — anything bigger than
# 20MB-ish would need File API rework. AgentBuff's 25MB cap is in the
# zone where inline still works.


_VIDEO_MIME_MAP = {
    "video/mp4": "video/mp4",
    "video/mpeg": "video/mpeg",
    "video/quicktime": "video/quicktime",
    "video/mov": "video/quicktime",
    "video/webm": "video/webm",
    "video/x-msvideo": "video/mp4",
    "video/avi": "video/mp4",
}


def _detect_video_mime_from_path(file_path: str) -> str:
    ext = Path(file_path).suffix.lower()
    mapping = {
        ".mp4": "video/mp4",
        ".mpeg": "video/mpeg", ".mpg": "video/mpeg",
        ".mov": "video/quicktime", ".qt": "video/quicktime",
        ".webm": "video/webm",
        ".avi": "video/mp4",
        ".mkv": "video/mp4",
    }
    return mapping.get(ext, "video/mp4")


def _video_default_prompt(user_prompt: Optional[str]) -> str:
    if user_prompt and user_prompt.strip():
        return user_prompt.strip()
    return (
        "Describe this video in detail. Note key moments, on-screen text, "
        "spoken audio if any, and overall content. Be thorough but concise."
    )


def _describe_video_openai_compat(
    provider_id: str,
    video_bytes: bytes,
    mime: str,
    prompt: str,
) -> tuple[Optional[str], Optional[str]]:
    """Generic chat.completions with `type: "video_url"` part.

    Covers Gemini (via openai-compatible mode), Qwen, Moonshot. For
    Gemini we use the generateContent native shape instead (separate
    function below) because video_url part isn't documented for the
    google generative endpoint — only the chat-completions adapter
    accepts it.
    """
    api_key = _get_provider_api_key(provider_id)
    if not api_key:
        return (None, "no_api_key")
    base_url = (
        os.environ.get(f"{provider_id.upper()}_BASE_URL")
        or PROVIDER_BASE_URLS.get(provider_id)
    )
    if not base_url:
        return (None, "no base_url")
    model = (
        os.environ.get(f"{provider_id.upper()}_VIDEO_MODEL")
        or DEFAULT_MODELS_BY_CAPABILITY["video"].get(provider_id)
    )
    if not model:
        return (None, "no default video model")
    canon_mime = _VIDEO_MIME_MAP.get((mime or "").lower(), "video/mp4")
    try:
        b64 = base64.b64encode(video_bytes).decode("ascii")
        body = {
            "model": model,
            "messages": [
                {
                    "role": "user",
                    "content": [
                        {"type": "text", "text": prompt},
                        {
                            "type": "video_url",
                            "video_url": {
                                "url": f"data:{canon_mime};base64,{b64}",
                            },
                        },
                    ],
                }
            ],
            "max_tokens": 1024,
        }
        resp = _http_post(
            f"{base_url.rstrip('/')}/chat/completions",
            provider_id=provider_id,
            headers={
                "Authorization": f"Bearer {api_key}",
                "Content-Type": "application/json",
            },
            json=body,
        )
        if resp.status_code != 200:
            return (None, f"HTTP {resp.status_code}: {resp.text[:200]}")
        data = resp.json()
        choices = data.get("choices") or []
        if not choices:
            return (None, "no choices")
        message = choices[0].get("message") or {}
        content = message.get("content")
        if isinstance(content, str):
            return (content.strip(), None)
        if isinstance(content, list):
            text = "\n".join(
                (p.get("text") or "").strip()
                for p in content
                if isinstance(p, dict) and p.get("text")
            ).strip()
            if text:
                return (text, None)
        reasoning = message.get("reasoning_content")
        if isinstance(reasoning, str) and reasoning.strip():
            return (reasoning.strip(), None)
        return (None, "empty content")
    except Exception as exc:
        return (None, f"{type(exc).__name__}: {exc}")


def _describe_gemini_video(
    video_bytes: bytes, mime: str, prompt: str
) -> tuple[Optional[str], Optional[str]]:
    """Gemini generateContent with inline_data for video. Same shape as
    image, just larger MIME family."""
    api_key = _get_provider_api_key("gemini")
    if not api_key:
        return (None, "no_api_key")
    model = (
        os.environ.get("GEMINI_VIDEO_MODEL")
        or DEFAULT_MODELS_BY_CAPABILITY["video"]["gemini"]
    )
    canon_mime = _VIDEO_MIME_MAP.get((mime or "").lower(), "video/mp4")
    try:
        b64 = base64.b64encode(video_bytes).decode("ascii")
        body = {
            "contents": [
                {
                    "parts": [
                        {"text": prompt},
                        {
                            "inline_data": {
                                "mime_type": canon_mime,
                                "data": b64,
                            }
                        },
                    ]
                }
            ],
            "generationConfig": {"temperature": 0.1, "maxOutputTokens": 1024},
        }
        url = (
            f"{PROVIDER_BASE_URLS['gemini']}/models/{model}:generateContent"
            f"?key={api_key}"
        )
        resp = _http_post(url, provider_id="gemini", json=body)
        if resp.status_code != 200:
            return (None, f"HTTP {resp.status_code}: {resp.text[:200]}")
        data = resp.json()
        candidates = data.get("candidates") or []
        if not candidates:
            return (None, "no candidates")
        parts = (candidates[0].get("content") or {}).get("parts") or []
        text = "".join(
            (p.get("text") or "") for p in parts if isinstance(p, dict)
        ).strip()
        return ((text or None), None if text else "empty content")
    except Exception as exc:
        return (None, f"{type(exc).__name__}: {exc}")


PROVIDER_DESCRIBE_VIDEO_FNS: dict[
    str,
    Callable[[bytes, str, str], tuple[Optional[str], Optional[str]]],
] = {
    "gemini": _describe_gemini_video,
    "google": _describe_gemini_video,
    "qwen": lambda b, m, p: _describe_video_openai_compat("qwen", b, m, p),
    "moonshot": lambda b, m, p: _describe_video_openai_compat(
        "moonshot", b, m, p
    ),
}


def _run_video_chain(
    file_path: str, prompt: str
) -> tuple[Optional[str], Optional[str], list[str]]:
    if not _capability_enabled("video"):
        return (None, None, ["video capability disabled in config"])

    video_bytes = _read_file_cached(file_path)
    if video_bytes is None:
        return (None, None, [f"read_file: failed to read {file_path}"])
    mime = _detect_video_mime_from_path(file_path)
    attempts: list[str] = []
    tried: set[str] = set()

    def _canon(pid: str) -> str:
        return "gemini" if pid == "google" else pid

    def _try(provider_id: str) -> tuple[Optional[str], Optional[str]]:
        canonical = _canon(provider_id)
        if canonical in tried:
            return (None, "already_tried")
        tried.add(canonical)
        fn = PROVIDER_DESCRIBE_VIDEO_FNS.get(provider_id)
        if not fn:
            attempts.append(f"{provider_id}: not registered")
            return (None, "not_registered")
        text, err = fn(video_bytes, mime, prompt)
        if text:
            attempts.append(f"{provider_id}: ok ({len(text)} chars)")
        else:
            attempts.append(f"{provider_id}: {err or 'empty'}")
        return (text, err)

    active = _get_active_chat_provider()
    if active:
        text, _ = _try(active)
        if text:
            return (text, active, attempts)

    priority = AUTO_PRIORITY_BY_CAPABILITY["video"]
    items = sorted(priority.items(), key=lambda kv: (kv[1], kv[0]))
    for provider_id, _prio in items:
        text, _ = _try(provider_id)
        if text:
            return (text, provider_id, attempts)

    return (None, None, attempts)


# ══════════════════════════════════════════════════════════════════════
# DOCUMENT chain — PDF/DOCX/XLSX/PPTX
# ══════════════════════════════════════════════════════════════════════
#
# Two strategies, depending on active chat model:
#
#   STRATEGY A — NATIVE PASSTHROUGH (Anthropic Claude, Gemini)
#     The model accepts PDF inline via its multimodal API; we don't
#     need to extract text. Hermes' built-in `context_note` pattern
#     already injects "[document at PATH]" and the agent has tools
#     that can attach files. We just signal to Hermes which path to
#     take by leaving the message alone.
#
#   STRATEGY B — TEXT EXTRACTION (everyone else)
#     For non-PDF-native models (OpenAI, DeepSeek, Mistral, Groq,
#     local, etc.), we extract text via Python libraries and inline
#     the extracted text into the message text. The agent then sees
#     the document content directly.
#
# Libraries used (all optional — if not installed, fall through):
#   * pdfplumber — PDF (preferred; better layout preservation)
#   * pypdf      — PDF (fallback if pdfplumber missing)
#   * docx       — python-docx, .docx files
#   * openpyxl   — .xlsx files
#   * pptx       — python-pptx, .pptx files
#
# Cap output at MAX_DOC_EXTRACT_CHARS to avoid blowing the context
# window. OpenClaw's behaviour was similar (~50k char default).


MAX_DOC_EXTRACT_CHARS = 50_000

_DOCUMENT_MIME_HINTS = {
    "application/pdf": "pdf",
    "application/x-pdf": "pdf",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "docx",
    "application/msword": "docx",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": "xlsx",
    "application/vnd.ms-excel": "xlsx",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
    "application/vnd.ms-powerpoint": "pptx",
}

_DOCUMENT_EXT_HINTS = {
    ".pdf": "pdf",
    ".docx": "docx",
    ".doc": "docx",
    ".xlsx": "xlsx",
    ".xls": "xlsx",
    ".pptx": "pptx",
    ".ppt": "pptx",
}


def _classify_document(file_path: str, mime: str) -> Optional[str]:
    """Map a (path, mime) pair to one of: pdf, docx, xlsx, pptx, or None."""
    canon = (mime or "").lower()
    kind = _DOCUMENT_MIME_HINTS.get(canon)
    if kind:
        return kind
    return _DOCUMENT_EXT_HINTS.get(Path(file_path).suffix.lower())


def _extract_pdf_text(file_path: str) -> tuple[Optional[str], Optional[str]]:
    """Extract text from PDF via pdfplumber → pypdf fallback."""
    # Try pdfplumber first (better layout handling)
    try:
        import pdfplumber  # type: ignore[import]
        try:
            with pdfplumber.open(file_path) as pdf:
                pages_text: list[str] = []
                for i, page in enumerate(pdf.pages):
                    try:
                        t = page.extract_text() or ""
                    except Exception:
                        t = ""
                    if t.strip():
                        pages_text.append(f"--- Page {i + 1} ---\n{t.strip()}")
                    if sum(len(p) for p in pages_text) >= MAX_DOC_EXTRACT_CHARS:
                        break
                combined = "\n\n".join(pages_text).strip()
                if combined:
                    return (combined[:MAX_DOC_EXTRACT_CHARS], None)
                return (None, "pdf had no extractable text")
        except Exception as exc:
            return (None, f"pdfplumber: {type(exc).__name__}: {exc}")
    except ImportError:
        pass

    # Fallback: pypdf
    try:
        import pypdf  # type: ignore[import]
        try:
            reader = pypdf.PdfReader(file_path)
            pages_text = []
            for i, page in enumerate(reader.pages):
                try:
                    t = page.extract_text() or ""
                except Exception:
                    t = ""
                if t.strip():
                    pages_text.append(f"--- Page {i + 1} ---\n{t.strip()}")
                if sum(len(p) for p in pages_text) >= MAX_DOC_EXTRACT_CHARS:
                    break
            combined = "\n\n".join(pages_text).strip()
            if combined:
                return (combined[:MAX_DOC_EXTRACT_CHARS], None)
            return (None, "pdf had no extractable text")
        except Exception as exc:
            return (None, f"pypdf: {type(exc).__name__}: {exc}")
    except ImportError:
        return (None, "no PDF library available (pdfplumber/pypdf)")


def _extract_docx_text(file_path: str) -> tuple[Optional[str], Optional[str]]:
    try:
        import docx  # type: ignore[import]
    except ImportError:
        return (None, "python-docx not installed")
    try:
        doc = docx.Document(file_path)
        chunks = [p.text for p in doc.paragraphs if p.text.strip()]
        # Also include table cells
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        chunks.append(cell.text.strip())
        combined = "\n".join(chunks).strip()
        if not combined:
            return (None, "docx had no extractable text")
        return (combined[:MAX_DOC_EXTRACT_CHARS], None)
    except Exception as exc:
        return (None, f"docx: {type(exc).__name__}: {exc}")


def _extract_xlsx_text(file_path: str) -> tuple[Optional[str], Optional[str]]:
    try:
        import openpyxl  # type: ignore[import]
    except ImportError:
        return (None, "openpyxl not installed")
    try:
        wb = openpyxl.load_workbook(file_path, data_only=True, read_only=True)
        chunks: list[str] = []
        total = 0
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            chunks.append(f"--- Sheet: {sheet_name} ---")
            for row in ws.iter_rows(values_only=True):
                row_text = " | ".join(
                    "" if v is None else str(v) for v in row
                ).strip(" |")
                if row_text:
                    chunks.append(row_text)
                    total += len(row_text)
                    if total >= MAX_DOC_EXTRACT_CHARS:
                        chunks.append("...(truncated, file too large)")
                        wb.close()
                        return (
                            "\n".join(chunks)[:MAX_DOC_EXTRACT_CHARS],
                            None,
                        )
            chunks.append("")
        wb.close()
        combined = "\n".join(chunks).strip()
        if not combined:
            return (None, "xlsx had no extractable text")
        return (combined[:MAX_DOC_EXTRACT_CHARS], None)
    except Exception as exc:
        return (None, f"openpyxl: {type(exc).__name__}: {exc}")


def _extract_pptx_text(file_path: str) -> tuple[Optional[str], Optional[str]]:
    try:
        from pptx import Presentation  # type: ignore[import]
    except ImportError:
        return (None, "python-pptx not installed")
    try:
        prs = Presentation(file_path)
        chunks: list[str] = []
        for i, slide in enumerate(prs.slides):
            chunks.append(f"--- Slide {i + 1} ---")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    chunks.append(shape.text.strip())
            chunks.append("")
        combined = "\n".join(chunks).strip()
        if not combined:
            return (None, "pptx had no extractable text")
        return (combined[:MAX_DOC_EXTRACT_CHARS], None)
    except Exception as exc:
        return (None, f"pptx: {type(exc).__name__}: {exc}")


_DOC_EXTRACTORS: dict[
    str, Callable[[str], tuple[Optional[str], Optional[str]]]
] = {
    "pdf": _extract_pdf_text,
    "docx": _extract_docx_text,
    "xlsx": _extract_xlsx_text,
    "pptx": _extract_pptx_text,
}


def _active_supports_native_pdf() -> bool:
    """True if `model.default` provider can ingest PDF inline (no
    extraction needed)."""
    active = _get_active_chat_provider()
    return active in {"anthropic", "gemini", "google"}


def _run_document_extraction(
    file_path: str, mime: str
) -> tuple[Optional[str], Optional[str], Optional[str]]:
    """Extract document text. Returns (text, doc_kind, error).

    text=None means extraction failed OR no library available — caller
    falls back to default "agent uses file tool" path.
    """
    if not _capability_enabled("document"):
        return (None, None, "document capability disabled in config")
    doc_kind = _classify_document(file_path, mime)
    if doc_kind is None:
        return (None, None, "not a recognized document type")
    extractor = _DOC_EXTRACTORS.get(doc_kind)
    if extractor is None:
        return (None, doc_kind, f"no extractor for {doc_kind}")
    text, err = extractor(file_path)
    return (text, doc_kind, err)


# ══════════════════════════════════════════════════════════════════════
# Monkey-patch installers
# ══════════════════════════════════════════════════════════════════════


def _make_patched_transcribe_audio(original):
    """Wrap Hermes' transcribe_audio with our extended chain fallback."""
    def patched(file_path: str, model: Optional[str] = None):
        # 0. Binary cache lookup — if preflight (hook) already
        #    transcribed this exact file, return that result. Avoids
        #    double API call.
        cached = _BINARY_CACHE.get_transcript(file_path)
        if cached is not None and cached.get("success"):
            logger.debug(
                "agentbuff-multimodal: transcribe_audio cache hit for %s",
                file_path,
            )
            return cached

        # 1. Try Hermes' original first. If the user has explicitly
        #    configured `stt.provider` and the configured provider works,
        #    that path will succeed. We add value when original returns
        #    "no provider configured" / "stt disabled" / network errors.
        original_result = None
        try:
            original_result = original(file_path, model=model)
        except TypeError:
            # Maybe the signature changed in a Hermes upgrade. Fall back
            # to positional call with just file_path.
            try:
                original_result = original(file_path)
            except Exception as exc:
                logger.debug(
                    "agentbuff-multimodal: original transcribe_audio(positional) "
                    "raised %s — running extended chain",
                    exc,
                )
        except Exception as exc:
            logger.debug(
                "agentbuff-multimodal: original transcribe_audio raised %s — "
                "running extended chain",
                exc,
            )

        if isinstance(original_result, dict) and original_result.get("success"):
            return original_result

        # 2. Hermes' chain didn't produce a transcript. Try our extended
        #    chain (Gemini/Deepgram/etc + active-chat-model fallback).
        try:
            transcript, provider, attempts = _run_extended_chain(file_path)
        except Exception as exc:
            logger.warning(
                "agentbuff-multimodal: extended chain crashed: %s",
                exc, exc_info=True,
            )
            transcript, provider, attempts = (None, None, [f"crash: {exc}"])

        if transcript:
            logger.info(
                "agentbuff-multimodal: transcribed via %s (%d chars)",
                provider, len(transcript),
            )
            result_dict = {
                "success": True,
                "transcript": transcript,
                "provider": f"agentbuff-{provider}",
                "agentbuff_attempts": attempts,
            }
            _BINARY_CACHE.put_transcript(file_path, result_dict)
            return result_dict
        if transcript == "":
            # Silent audio — propagate to caller as success with empty text
            return {
                "success": True,
                "transcript": "",
                "provider": f"agentbuff-{provider or 'silent'}",
            }

        # 3. Both chains failed. Return Hermes' original result if we have
        #    one (preserves its error message), else a synthesized failure
        #    that aggregates our attempts log for diagnostics.
        if isinstance(original_result, dict):
            return original_result
        diag = " | ".join(attempts[-4:]) if attempts else "no providers"
        return {
            "success": False,
            "transcript": "",
            "error": f"All STT providers failed: {diag}",
            "agentbuff_attempts": attempts,
        }

    setattr(patched, _PATCHED_SENTINEL, True)
    return patched


def _make_patched_vision_analyze_tool(original):
    """Wrap Hermes' vision_analyze_tool (async) with our extended chain.

    Hermes contract (`tools/vision_tools.py:633`):
        async def vision_analyze_tool(image_url, user_prompt, model=None) -> str

    Returns a JSON string with shape `{"success": bool, "analysis": str,
    "error": str}`. The function accepts http(s)://, file://, or local
    file paths. We only run the extended chain when given a LOCAL path
    or `file://` URI — remote URLs need Hermes' download infra so we
    let original handle those.
    """
    import asyncio  # local import to avoid module load order issues
    import json as _json

    async def patched(
        image_url: str,
        user_prompt: str,
        model: Optional[str] = None,
    ) -> str:
        # 1. Try Hermes' original first.
        original_result = None
        try:
            original_result = await original(image_url, user_prompt, model=model)
        except TypeError:
            try:
                original_result = await original(image_url, user_prompt)
            except Exception as exc:
                logger.debug(
                    "agentbuff-multimodal: original vision_analyze_tool(positional) "
                    "raised %s — running extended chain",
                    exc,
                )
        except Exception as exc:
            logger.debug(
                "agentbuff-multimodal: original vision_analyze_tool raised %s — "
                "running extended chain",
                exc,
            )

        # Parse Hermes' JSON envelope to decide if we need the extended chain.
        if isinstance(original_result, str):
            try:
                parsed = _json.loads(original_result)
                if isinstance(parsed, dict) and parsed.get("success"):
                    return original_result
            except Exception:
                # Original returned non-JSON; treat as failure and continue
                pass

        # 2. Hermes' chain failed. Try our extended chain — but only for
        #    local paths. Remote URLs would need download infra we don't
        #    duplicate; the original already failed those.
        local_path = _resolve_local_image_path(image_url)
        if not local_path:
            if isinstance(original_result, str):
                return original_result
            return _json.dumps({
                "success": False,
                "analysis": "",
                "error": "vision_analyze: remote URL and original failed",
            })

        try:
            description, provider, attempts = await asyncio.to_thread(
                _run_vision_chain,
                local_path,
                _vision_default_prompt(user_prompt),
            )
        except Exception as exc:
            logger.warning(
                "agentbuff-multimodal: extended vision chain crashed: %s",
                exc, exc_info=True,
            )
            description, provider, attempts = (None, None, [f"crash: {exc}"])

        if description:
            logger.info(
                "agentbuff-multimodal: described image via %s (%d chars)",
                provider, len(description),
            )
            return _json.dumps({
                "success": True,
                "analysis": description,
                "provider": f"agentbuff-{provider}",
                "agentbuff_attempts": attempts,
            })

        # Both chains failed. Prefer original's message if present.
        if isinstance(original_result, str):
            return original_result
        diag = " | ".join(attempts[-4:]) if attempts else "no providers"
        return _json.dumps({
            "success": False,
            "analysis": "",
            "error": f"All vision providers failed: {diag}",
            "agentbuff_attempts": attempts,
        })

    setattr(patched, _PATCHED_SENTINEL, True)
    return patched


def _resolve_local_image_path(image_url: str) -> Optional[str]:
    """Convert a `file://` URI or local path to a real filesystem path.
    Returns None for http(s) URLs (let Hermes handle remote)."""
    if not isinstance(image_url, str) or not image_url:
        return None
    if image_url.startswith(("http://", "https://")):
        return None
    if image_url.startswith("file://"):
        path = image_url[7:]
        # Strip Windows-style leading slash if present (file:///C:/...)
        if len(path) > 2 and path[0] == "/" and path[2] == ":":
            path = path[1:]
        return path if Path(path).is_file() else None
    p = Path(image_url).expanduser()
    return str(p) if p.is_file() else None


def _install_patches() -> None:
    """Patch Hermes' transcription + vision functions. Idempotent."""
    _patch_function(
        module_name="tools.transcription_tools",
        attr_name="transcribe_audio",
        wrapper_factory=_make_patched_transcribe_audio,
        label="STT",
    )
    _patch_function(
        module_name="tools.vision_tools",
        attr_name="vision_analyze_tool",
        wrapper_factory=_make_patched_vision_analyze_tool,
        label="vision",
    )
    # gTTS provider extension — primary TTS provider after edge-tts
    # broke at the library layer (Microsoft auth change 2026-05-23
    # returns 403 on every handshake). Lives in sibling file
    # `tts_gtts.py` so it's self-contained + easy to remove when
    # upstream edge-tts patches the auth.
    try:
        from . import tts_gtts  # type: ignore
        tts_gtts.install_gtts_patch()
    except Exception:
        logger.exception(
            "agentbuff-multimodal: tts_gtts.install_gtts_patch() crashed "
            "— TTS will fall back to Hermes' built-in providers (edge/openai/etc)",
        )
    # Interactive prompts (clarify) — make Hermes emit `clarify.request`
    # wire event so /app web UI can render an interactive button card.
    # Without this, Hermes' clarify_tool silently returns "" for
    # operator-mode sessions because no platform adapter is registered
    # (verified live: only `approval.request` emits on the wire by
    # default — see `tui_gateway/server.py:1944`). Lives in sibling
    # file `interactive_bridge.py`.
    try:
        from . import interactive_bridge  # type: ignore
        interactive_bridge.install_interactive_patches()
    except Exception:
        logger.exception(
            "agentbuff-multimodal: interactive_bridge.install_interactive_patches() "
            "crashed — clarify prompts silently no-op in /app",
        )


def _patch_function(
    *,
    module_name: str,
    attr_name: str,
    wrapper_factory: Callable[[Any], Any],
    label: str,
) -> None:
    """Generic monkey-patch installer: replace `module.attr` with
    `wrapper_factory(original)` and rebind top-level importers."""
    try:
        module = __import__(module_name, fromlist=[attr_name])
    except ImportError as exc:
        logger.warning(
            "agentbuff-multimodal: cannot import %s (%s) — %s patch skipped",
            module_name, exc, label,
        )
        return

    original = getattr(module, attr_name, None)
    if not callable(original):
        logger.warning(
            "agentbuff-multimodal: %s has no callable %s — %s patch skipped",
            module_name, attr_name, label,
        )
        return
    if getattr(original, _PATCHED_SENTINEL, False):
        logger.debug(
            "agentbuff-multimodal: %s already patched, skipping",
            attr_name,
        )
        return

    patched = wrapper_factory(original)
    setattr(module, attr_name, patched)

    # Rebind top-level importers (modules that did
    # `from X import Y` at top level).
    rebind_count = 0
    for mod_name, mod in list(sys.modules.items()):
        if mod is None or mod is module:
            continue
        try:
            if getattr(mod, attr_name, None) is original:
                setattr(mod, attr_name, patched)
                rebind_count += 1
        except Exception:
            continue

    logger.info(
        "agentbuff-multimodal: patched %s.%s (rebound %d top-level importer(s))",
        module_name, attr_name, rebind_count,
    )


# ══════════════════════════════════════════════════════════════════════
# pre_gateway_dispatch hook — auto-enrich video + document for channels
# ══════════════════════════════════════════════════════════════════════
#
# Audio + image are auto-enriched by Hermes itself (gateway calls
# `transcribe_audio` / `vision_analyze_tool` BEFORE the agent sees the
# message; our monkey-patch extends the providers used there).
#
# Video + document are NOT auto-enriched by Hermes — the gateway just
# caches them and prepends a path note. To make these "just work" for
# mass-market users (chief's "carry the task" promise), we run our own
# chains inside the `pre_gateway_dispatch` hook and rewrite the event
# text with the description / extracted content.
#
# Hook contract per `gateway/run.py:6439-6478`:
#     def _on_pre_gateway_dispatch(event, gateway, session_store, **kw)
#         → None  /  {"action": "skip"|"rewrite"|"allow", ...}


def _on_pre_gateway_dispatch(
    event: Any = None,
    gateway: Any = None,
    session_store: Any = None,
    **kwargs: Any,
) -> Optional[dict]:
    """Pre-dispatch enrichment for video + document + (optionally) audio.

    Three behaviours, all gated by config:

      1. Video + document (always-on) — Hermes doesn't auto-enrich
         these, so we run our own chains and prepend context notes.

      2. Audio preflight (off by default; enable with
         `tools.media.audio.preflight: true`) — transcribe audio
         INLINE here so the transcript is in `event.text` before
         Hermes' auth/mention check runs. Useful for group chats that
         require mention (audio with no caption text would fail mention
         check without preflight).

      3. Echo transcript (off by default; enable with
         `tools.media.audio.echoTranscript: true`) — send a separate
         `📝 "..."` confirmation message to the user's channel so they
         see what the bot heard. Fire-and-forget via
         `asyncio.create_task(adapter.send(...))`.

    All three run in parallel when multiple attachments are present
    (concurrent.futures.ThreadPoolExecutor, capped by MEDIA_CONCURRENCY).

    Runs in plugin context. Never raises — any failure logs and returns
    None (let Hermes' default flow run).
    """
    try:
        media_urls = list(getattr(event, "media_urls", []) or [])
        media_types = list(getattr(event, "media_types", []) or [])
        original_text = getattr(event, "text", "") or ""

        if not media_urls:
            return None

        # Decide per-attachment whether to process in hook OR delegate
        # to Hermes' built-in enrichment (our monkey-patches kick in
        # there for audio + image).
        audio_preflight = _capability_audio_preflight()
        echo_enabled = _capability_echo_transcript()

        tasks_by_idx: list[tuple[int, str, Callable[[], Optional[str]], bool]] = []
        # tuple = (idx, kind, callable that returns prefix-note-or-None,
        #          should_emit_echo)

        for idx, path in enumerate(media_urls):
            mime = (
                media_types[idx]
                if idx < len(media_types) and media_types[idx]
                else ""
            )
            mime_lower = mime.lower()

            if mime_lower.startswith("audio/"):
                # Process in hook ONLY if preflight or echo is enabled.
                # Otherwise let Hermes' enrich + our monkey-patched
                # transcribe_audio handle it.
                if audio_preflight or echo_enabled:
                    tasks_by_idx.append(
                        (idx, "audio",
                         lambda p=path: _enrich_audio_inline(p),
                         echo_enabled)
                    )
                continue

            if mime_lower.startswith("image/"):
                # Delegated to vision_analyze_tool monkey-patch (always).
                continue

            if mime_lower.startswith("video/"):
                tasks_by_idx.append(
                    (idx, "video",
                     lambda p=path: _enrich_video(p),
                     False)
                )
                continue

            # Document path — non-PDF or non-native-PDF model
            tasks_by_idx.append(
                (idx, "document",
                 lambda p=path, m=mime_lower: _enrich_document(p, m),
                 False)
            )

        if not tasks_by_idx:
            return None

        # Concurrent execution — fan out provider calls in parallel
        # when multiple attachments are in one event. Bounded by
        # MEDIA_CONCURRENCY env (default 3).
        callables = [t[2] for t in tasks_by_idx]
        results = _run_concurrently(callables)

        prefix_parts: list[str] = []
        echo_lines: list[tuple[str, str]] = []  # (path, transcript)
        for (idx, kind, _, should_echo), result in zip(tasks_by_idx, results):
            if isinstance(result, Exception):
                logger.warning(
                    "agentbuff-multimodal: hook task %d (%s) raised: %s",
                    idx, kind, result,
                )
                continue
            if isinstance(result, str) and result:
                prefix_parts.append(result)
                # If this was an audio preflight + echo enabled, extract
                # the transcript from the note to send as separate echo
                # message via the channel adapter.
                if should_echo and kind == "audio":
                    transcript = _extract_transcript_from_note(result)
                    if transcript:
                        echo_lines.append((media_urls[idx], transcript))

        # Fire-and-forget echo transcripts via the channel adapter
        if echo_lines and gateway is not None and event is not None:
            _schedule_echo_messages(gateway, event, echo_lines)

        if not prefix_parts:
            return None

        new_text = "\n\n".join(prefix_parts)
        if original_text.strip():
            new_text = f"{new_text}\n\n{original_text}"

        logger.info(
            "agentbuff-multimodal: pre_gateway_dispatch enriched %d "
            "attachment(s) (added %d chars to message; %d echo)",
            len(prefix_parts), len(new_text) - len(original_text),
            len(echo_lines),
        )
        return {"action": "rewrite", "text": new_text}
    except Exception as exc:
        logger.warning(
            "agentbuff-multimodal: pre_gateway_dispatch hook crashed: %s",
            exc, exc_info=True,
        )
        return None


def _enrich_audio_inline(path: str) -> Optional[str]:
    """Transcribe audio in the hook (preflight) and return a channel-
    style context note. Stores result in cache so Hermes' enrich call
    later sees the cached value and doesn't re-transcribe.
    """
    try:
        transcript, provider, attempts = _run_extended_chain(path)
    except Exception as exc:
        logger.warning(
            "agentbuff-multimodal: audio preflight crashed for %s: %s",
            path, exc,
        )
        return None

    if transcript:
        result_dict = {
            "success": True,
            "transcript": transcript,
            "provider": f"agentbuff-{provider}",
            "agentbuff_attempts": attempts,
        }
        _BINARY_CACHE.put_transcript(path, result_dict)
        logger.info(
            "agentbuff-multimodal: preflight-transcribed %s via %s "
            "(%d chars)",
            path, provider, len(transcript),
        )
        return (
            f'[The user sent a voice message~ Here\'s what they said: '
            f'"{transcript}"]'
        )
    return None


def _extract_transcript_from_note(note: str) -> Optional[str]:
    """Parse the channel-style note `[... Here's what they said: "X"]`
    to extract just X for use in echo message."""
    # Match pattern: Here's what they said: "..."
    m = re.search(r'they said:\s*"([^"]*)"', note)
    if m:
        return m.group(1).strip()
    return None


def _schedule_echo_messages(
    gateway: Any, event: Any, echo_lines: list[tuple[str, str]]
) -> None:
    """Schedule fire-and-forget echo messages to the channel.
    Format: `📝 "transcript"` — matches OpenClaw's
    `DEFAULT_ECHO_TRANSCRIPT_FORMAT` constant.

    Uses gateway.adapters[platform].send(chat_id, content) which is
    Hermes' standard outbound API. If anything fails (no adapter, no
    chat_id, send error), logged but never raises.
    """
    try:
        import asyncio
        source = getattr(event, "source", None)
        if source is None:
            return
        platform = getattr(source, "platform", None)
        chat_id = getattr(source, "chat_id", None)
        if platform is None or chat_id is None:
            return
        adapters = getattr(gateway, "adapters", None)
        if not adapters:
            return
        adapter = adapters.get(platform)
        if adapter is None or not hasattr(adapter, "send"):
            return

        try:
            loop = asyncio.get_running_loop()
        except RuntimeError:
            # No running loop — can't schedule. Skip silently.
            return

        for _path, transcript in echo_lines:
            content = f'📝 "{transcript}"'
            try:
                loop.create_task(
                    adapter.send(chat_id=chat_id, content=content)
                )
            except Exception as exc:
                logger.debug(
                    "agentbuff-multimodal: failed to schedule echo: %s", exc
                )
    except Exception as exc:
        logger.debug(
            "agentbuff-multimodal: _schedule_echo_messages outer crash: %s",
            exc,
        )


def _enrich_video(path: str) -> Optional[str]:
    """Run video chain, return channel-style context note. None on failure."""
    try:
        description, provider, attempts = _run_video_chain(
            path, _video_default_prompt(None)
        )
    except Exception as exc:
        logger.warning(
            "agentbuff-multimodal: video enrichment crashed for %s: %s",
            path, exc,
        )
        return None

    if description:
        logger.info(
            "agentbuff-multimodal: video described via %s (%d chars)",
            provider, len(description),
        )
        return (
            f'[The user sent a video. Here\'s what\'s in it: "{description}"]'
        )
    diag = " | ".join(attempts[-3:]) if attempts else "no providers"
    return (
        f"[The user sent a video at {path}. AgentBuff couldn't auto-describe "
        f"it ({diag}). Ask the user what they want you to do with it, or use "
        f"the video_analyze tool manually.]"
    )


def _enrich_document(path: str, mime: str) -> Optional[str]:
    """Document enrichment with two strategies:

      A) Native PDF passthrough — if active chat model is anthropic
         claude or gemini, AND the document is PDF, leave Hermes' default
         flow alone (it will inject a "[document at PATH]" note and the
         model reads PDF natively via its own multimodal API).

      B) Text extraction — for any other model (or non-PDF docs), extract
         text via pdfplumber/python-docx/openpyxl/python-pptx and inline
         the extracted text into the message.
    """
    doc_kind = _classify_document(path, mime)
    if doc_kind is None:
        return None  # Not a recognized document type — let Hermes handle

    # Strategy A: native passthrough for PDF + native model
    if doc_kind == "pdf" and _active_supports_native_pdf():
        # Return None so Hermes' default path runs (which adds context
        # note pointing at file). The agent's active model will read the
        # PDF directly via its multimodal API once it pulls the file.
        return None

    # Strategy B: extract text
    text, _kind, err = _run_document_extraction(path, mime)
    name = Path(path).name
    if text:
        logger.info(
            "agentbuff-multimodal: extracted %d chars from %s (%s)",
            len(text), name, doc_kind,
        )
        return (
            f"[The user sent a document: '{name}' ({doc_kind.upper()}). "
            f"Extracted content below — original file at {path}.]\n\n"
            f"--- BEGIN {name} ---\n{text}\n--- END {name} ---"
        )

    # Extraction failed — defer to Hermes' default note path
    logger.info(
        "agentbuff-multimodal: extraction skipped for %s (%s): %s",
        name, doc_kind, err,
    )
    return None


# ══════════════════════════════════════════════════════════════════════
# Plugin entry — register() called by Hermes' plugin loader
# ══════════════════════════════════════════════════════════════════════


def register(ctx: Any) -> None:
    """Hermes plugin entry point (called once during `discover_and_load`).

    The monkey-patches are already installed at import time (see bottom
    of this file). Here we just wire the pre_gateway_dispatch hook for
    video + document enrichment.

    ctx is a `PluginContext` (`hermes_cli/plugins.py:701-716`) that
    exposes `register_hook(name, callback)`.
    """
    try:
        ctx.register_hook("pre_gateway_dispatch", _on_pre_gateway_dispatch)
        logger.info(
            "agentbuff-multimodal: registered pre_gateway_dispatch hook "
            "(video + document auto-enrichment active)"
        )
    except Exception as exc:
        logger.exception(
            "agentbuff-multimodal: register_hook failed: %s — video/document "
            "auto-enrichment will NOT run, but audio + image patches stay active",
            exc,
        )


# Install patches at import time. Plugin loader runs this __init__.py
# once during `discover_and_load`. After that, every call across
# Hermes — channels + /app + agent-initiated tool calls — runs through
# our wrappers.
try:
    _install_patches()
except Exception as _exc:
    # Defensive: a plugin __init__ that raises would prevent the entire
    # plugin manager from advancing. Log and continue with original
    # behavior so we never break a working Hermes deployment.
    logger.exception(
        "agentbuff-multimodal: install_patches raised %s — extended chains not active",
        _exc,
    )


__all__ = [
    "DEFAULT_MODELS",
    "DEFAULT_MODELS_BY_CAPABILITY",
    "AUTO_PRIORITY",
    "AUTO_PRIORITY_BY_CAPABILITY",
    "PROVIDER_BASE_URLS",
    "PROVIDER_KEY_ENV_VARS",
    "register",
]
